#!/usr/bin/env python3
"""Generate a single executive slide for OAC/Essbase → Fabric migration services offering.

Usage:
    python scripts/generate_services_slide.py [--output path/to/slide.pptx]
"""
from __future__ import annotations

import argparse
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colours ──
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x10, 0x7C, 0x10)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
TEAL = RGBColor(0x00, 0xB7, 0xC3)
PURPLE = RGBColor(0x6B, 0x2D, 0x8B)
RED = RGBColor(0xD1, 0x34, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF2, 0xF1)
GRAY = RGBColor(0x60, 0x60, 0x60)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _shape(slide, l, t, w, h, fill, st=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def _txt(slide, l, t, w, h, text, sz=18, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, name="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return tb


def _phase_card(slide, x, y, w, h, num, title, duration, deliverables, color):
    """Phase card with number circle, title, duration, and bullet deliverables."""
    card = _shape(slide, x, y, w, h, WHITE)
    card.line.color.rgb = color
    card.line.width = Pt(2)

    # Number circle
    circle = _shape(slide, x + w / 2 - Inches(0.18), y + Inches(0.12),
                    Inches(0.36), Inches(0.36), color, MSO_SHAPE.OVAL)
    circle.text_frame.paragraphs[0].text = str(num)
    circle.text_frame.paragraphs[0].font.size = Pt(14)
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.name = "Segoe UI"
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    _txt(slide, x + Inches(0.08), y + Inches(0.52), w - Inches(0.16), Inches(0.3),
         title, sz=11, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Duration
    dur_bar = _shape(slide, x + Inches(0.15), y + Inches(0.82),
                     w - Inches(0.3), Inches(0.22), LIGHT)
    _txt(slide, x + Inches(0.15), y + Inches(0.82),
         w - Inches(0.3), Inches(0.22),
         duration, sz=8, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Deliverables
    dy = y + Inches(1.12)
    for d in deliverables:
        _txt(slide, x + Inches(0.12), dy, w - Inches(0.24), Inches(0.18),
             f"✓  {d}", sz=7.5, color=BLACK)
        dy += Inches(0.18)


def build(output: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)

    # ══ Title bar ══
    _shape(slide, 0, 0, SLIDE_W, Inches(1.05), DARK_BLUE, MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(0.6), Inches(0.12), Inches(9), Inches(0.5),
         "Migration Services — Oracle Analytics & Essbase → Microsoft Fabric",
         sz=22, bold=True, color=WHITE)
    _txt(slide, Inches(0.6), Inches(0.58), Inches(10), Inches(0.35),
         "End-to-end consulting engagement  —  AI-accelerated  —  Fixed-scope delivery",
         sz=12, color=TEAL)

    # ══ ROW 1: 5 engagement phases ══
    y1 = Inches(1.25)
    _txt(slide, Inches(0.4), y1, Inches(5), Inches(0.3),
         "Engagement Phases", sz=15, bold=True, color=DARK_BLUE)

    phases = [
        (1, "Discovery &\nAssessment", "2 weeks",
         ["OAC/Essbase inventory crawl",
          "Dependency & complexity map",
          "Migration readiness report",
          "Effort estimate & wave plan"],
         MS_BLUE),
        (2, "Architecture &\nDesign", "2 weeks",
         ["Fabric workspace design",
          "Data model & partition strategy",
          "Security model mapping",
          "ETL pipeline architecture"],
         TEAL),
        (3, "Migration &\nDevelopment", "8–12 weeks",
         ["Schema & data migration",
          "DAX semantic model build",
          "Report & dashboard conversion",
          "RLS/OLS security setup"],
         GREEN),
        (4, "Validation &\nUAT", "2–4 weeks",
         ["Data reconciliation (checksums)",
          "Visual regression testing",
          "Performance benchmarking",
          "User acceptance sign-off"],
         ORANGE),
        (5, "Go-Live &\nHypercare", "2–4 weeks",
         ["Production cutover",
          "Parallel run monitoring",
          "Knowledge transfer & training",
          "Post-go-live support"],
         PURPLE),
    ]
    px = Inches(0.4)
    pw = Inches(2.4)
    ph = Inches(2.8)
    for num, title, dur, dels, color in phases:
        _phase_card(slide, px, y1 + Inches(0.35), pw, ph,
                    num, title, dur, dels, color)
        # Connector arrow between cards
        if num < 5:
            _txt(slide, px + pw, y1 + Inches(1.2), Inches(0.2), Inches(0.4),
                 "›", sz=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        px += pw + Inches(0.18)

    # ══ ROW 2: Team & Deliverables + Value Proposition ══
    y2 = Inches(4.6)

    # Left: Team composition
    _txt(slide, Inches(0.4), y2, Inches(3.5), Inches(0.3),
         "Team Composition", sz=13, bold=True, color=DARK_BLUE)

    roles = [
        ("Engagement Lead", "Overall delivery, stakeholder management"),
        ("Data Architect", "Fabric workspace, Lakehouse & Warehouse design"),
        ("BI Developer", "Semantic models, DAX, Power BI reports"),
        ("ETL Engineer", "Data pipelines, PySpark notebooks"),
        ("QA / Validation", "Data reconciliation, regression testing"),
    ]
    ry = y2 + Inches(0.35)
    for role, desc in roles:
        bar = _shape(slide, Inches(0.4), ry, Inches(1.5), Inches(0.22),
                     MS_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
        bar.text_frame.paragraphs[0].text = f" {role}"
        bar.text_frame.paragraphs[0].font.size = Pt(7.5)
        bar.text_frame.paragraphs[0].font.bold = True
        bar.text_frame.paragraphs[0].font.color.rgb = WHITE
        bar.text_frame.paragraphs[0].font.name = "Segoe UI"
        _txt(slide, Inches(2.0), ry, Inches(2.2), Inches(0.22),
             desc, sz=7.5, color=GRAY)
        ry += Inches(0.26)

    # Center: Key deliverables
    _txt(slide, Inches(4.6), y2, Inches(3.5), Inches(0.3),
         "Key Deliverables", sz=13, bold=True, color=DARK_BLUE)

    deliverables = [
        ("Migration Readiness Report", "Asset inventory, complexity scores, wave plan"),
        ("Fabric Architecture Blueprint", "Workspace layout, security, data flow diagrams"),
        ("Migrated Artifacts", "Lakehouse tables, semantic models, PBI reports"),
        ("Validation Report", "Row counts, checksums, visual comparison"),
        ("Runbook & Training", "Operations guide, knowledge transfer sessions"),
    ]
    dy = y2 + Inches(0.35)
    for title, desc in deliverables:
        _txt(slide, Inches(4.6), dy, Inches(1.8), Inches(0.22),
             f"▸  {title}", sz=8, bold=True, color=DARK_BLUE)
        _txt(slide, Inches(6.4), dy, Inches(2.6), Inches(0.22),
             desc, sz=7.5, color=GRAY)
        dy += Inches(0.26)

    # Right: Value proposition
    _txt(slide, Inches(9.2), y2, Inches(3.8), Inches(0.3),
         "Why This Engagement", sz=13, bold=True, color=DARK_BLUE)

    values = [
        ("AI-Accelerated", "–50% migration time vs. manual approach", GREEN),
        ("Proven Tooling", "188+ modules, 4,005 automated tests", MS_BLUE),
        ("Risk-Managed", "Wave-based rollout with automated rollback", ORANGE),
        ("Full Coverage", "97% OAC objects, 185 mapping rules", TEAL),
        ("Knowledge Transfer", "Your team owns the outcome, not us", PURPLE),
    ]
    vy = y2 + Inches(0.35)
    for label, desc, color in values:
        dot = _shape(slide, Inches(9.2), vy + Inches(0.03),
                     Inches(0.14), Inches(0.14), color, MSO_SHAPE.OVAL)
        _txt(slide, Inches(9.45), vy, Inches(1.5), Inches(0.22),
             label, sz=8, bold=True, color=color)
        _txt(slide, Inches(10.9), vy, Inches(2.1), Inches(0.22),
             desc, sz=7.5, color=GRAY)
        vy += Inches(0.26)

    # ══ Bottom bar: CTA ══
    _shape(slide, 0, Inches(6.65), SLIDE_W, Inches(0.85), DARK_BLUE,
           MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(0.6), Inches(6.72), Inches(5), Inches(0.35),
         "Typical engagement: 16–24 weeks  •  3–6 consultants  •  Fixed-scope pricing available",
         sz=11, bold=True, color=WHITE)
    _txt(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
         "Start with a 2-week Discovery Workshop to assess your Oracle estate and build a migration roadmap.",
         sz=10, color=TEAL)

    # Right side CTA
    cta = _shape(slide, Inches(10.2), Inches(6.78), Inches(2.7), Inches(0.55),
                 TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
    cta.text_frame.paragraphs[0].text = "Schedule Discovery →"
    cta.text_frame.paragraphs[0].font.size = Pt(13)
    cta.text_frame.paragraphs[0].font.bold = True
    cta.text_frame.paragraphs[0].font.color.rgb = WHITE
    cta.text_frame.paragraphs[0].font.name = "Segoe UI"
    cta.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    prs.save(output)
    return output


def main():
    parser = argparse.ArgumentParser(
        description="Generate OAC→Fabric migration services one-slider")
    parser.add_argument("--output", "-o",
                        default="output/OAC_to_Fabric_Services.pptx")
    args = parser.parse_args()
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    path = build(args.output)
    print(f"✅ Services slide generated: {path}")
    print(f"   1 slide • widescreen 16:9 • {os.path.getsize(path):,} bytes")


if __name__ == "__main__":
    main()
