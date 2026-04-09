#!/usr/bin/env python3
"""Generate a Business Migration ROI PowerPoint deck.

Usage:
    python scripts/generate_business_deck.py [--output path/to/deck.pptx]

Produces an 11-slide executive deck covering:
  1. Title slide
  2. Migration accelerated with GenAI (project matrix + KPIs)
  3. What Can Be Migrated (object coverage)
  4. OAC & Essbase → Fabric Object Mapping Matrix
  5. Essbase & Smart View Migration
  6. Longview/EPM — Two Migration Paths (Option A vs B)
  7. Longview Migration Plans — Timeline & Savings
  8. Automation & Time Savings
  9. ROI & Business Impact
  10. Migration Platform Architecture
  11. Next Steps / Call to Action
"""
from __future__ import annotations

import argparse
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── Brand colours ──────────────────────────────────────────────────────
MICROSOFT_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_TEAL = RGBColor(0x00, 0xB7, 0xC3)
ACCENT_PURPLE = RGBColor(0x6B, 0x2D, 0x8B)
ACCENT_RED = RGBColor(0xD1, 0x34, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF2, 0xF1)
MEDIUM_GRAY = RGBColor(0x60, 0x60, 0x60)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────

def _set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color: RGBColor,
               shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_kpi_card(slide, left, top, width, height, title, value,
                  border_color, bg_color=WHITE):
    card = _add_shape(slide, left, top, width, height, bg_color)
    card.line.color.rgb = border_color
    card.line.width = Pt(3)

    # Title
    _add_textbox(slide, left + Inches(0.15), top + Inches(0.1),
                 width - Inches(0.3), Inches(0.4),
                 title, font_size=11, bold=True, color=border_color,
                 alignment=PP_ALIGN.CENTER)
    # Value
    _add_textbox(slide, left + Inches(0.1), top + Inches(0.45),
                 width - Inches(0.2), Inches(0.5),
                 value, font_size=22, bold=True, color=border_color,
                 alignment=PP_ALIGN.CENTER)


def _add_table(slide, left, top, width, rows_data, col_widths,
               header_color=MICROSOFT_BLUE, stripe_color=LIGHT_GRAY,
               row_height=0.4, font_size=12):
    """Add a styled table. rows_data: list of lists (first row = header)."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_height = Inches(row_height) * n_rows
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.text = str(cell_text)
            p.font.size = Pt(font_size)
            p.font.name = "Segoe UI"
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if ri == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                p.font.color.rgb = BLACK
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = stripe_color
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

    return shape


# ── Slide Builders ─────────────────────────────────────────────────────

def slide_title(prs: Presentation):
    """Slide 1: Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BLUE)

    _add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "OAC → Microsoft Fabric & Power BI",
                 font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "Business Migration Assessment — ROI & Capabilities",
                 font_size=24, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)

    # Tagline
    _add_textbox(slide, Inches(2), Inches(4.2), Inches(9), Inches(0.6),
                 "Fully automated, AI-accelerated migration  •  97% OAC object coverage  •  4,005 tests  •  185 object mappings",
                 font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Bottom bar
    _add_shape(slide, Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.7),
               MICROSOFT_BLUE, MSO_SHAPE.RECTANGLE)
    _add_textbox(slide, Inches(0.5), Inches(6.85), Inches(12), Inches(0.5),
                 "🚀 GitHub Copilot Accelerated   |   T-SQL → Spark SQL   •   PL/SQL → PySpark   •   "
                 "OAC RPD → Semantic Models   •   Data Quality   •   Auto Tests   •   Doc Generation",
                 font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_migration_matrix(prs: Presentation):
    """Slide 2: Migration accelerated with GenAI — project matrix + KPIs."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(8), Inches(0.6),
                 "Migration accelerated with GenAI",
                 font_size=30, bold=True, color=DARK_BLUE)

    _add_textbox(slide, Inches(8.5), Inches(0.25), Inches(4.5), Inches(0.5),
                 "Figures from other customers",
                 font_size=16, bold=True, color=ACCENT_ORANGE,
                 alignment=PP_ALIGN.RIGHT)

    # ── Project matrix table ──
    header = ["PROJECT", "TODAY", "TOMORROW", "W/o Copilot", "With Copilot", "Gain"]
    rows = [
        header,
        ["BI Finance", "Azure SQL", "Fabric Lakehouse", "5–7 wks", "3–4 wks", "–40%"],
        ["Enterprise DW\n(US + Brazil)", "Oracle Exadata → OAC", "Fabric Warehouse", "8–12 wks", "5–7 wks", "–40%"],
        ["Noram", "Oracle Analytics Cloud", "Fabric + Power BI", "5–7 wks", "3–4 wks", "–35%"],
        ["Informatica ETL", "Informatica", "Fabric Pipelines", "8–10 wks", "4–5 wks", "–50%"],
        ["Total Estimated", "", "", "8–12 months", "4–6 months", "~50%"],
    ]
    col_widths = [Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.3), Inches(1.3), Inches(0.9)]

    tbl_shape = _add_table(slide, Inches(0.3), Inches(1.0), Inches(8.3),
                           rows, col_widths)
    # Colour the Gain column cells
    table = tbl_shape.table
    for ri in range(1, len(rows)):
        gain_cell = table.cell(ri, 5)
        gain_cell.fill.solid()
        if ri == len(rows) - 1:
            gain_cell.fill.fore_color.rgb = ACCENT_GREEN
            # Total row green bg
            for ci in range(6):
                c = table.cell(ri, ci)
                c.fill.solid()
                c.fill.fore_color.rgb = ACCENT_GREEN
                c.text_frame.paragraphs[0].font.color.rgb = WHITE
                c.text_frame.paragraphs[0].font.bold = True
        else:
            gain_cell.fill.fore_color.rgb = ACCENT_GREEN
            gain_cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            gain_cell.text_frame.paragraphs[0].font.bold = True

    # ── KPI cards on the right ──
    kpi_data = [
        ("Infra TCO Reduction", "–50 to –65%", MICROSOFT_BLUE),
        ("Ops Effort Reduction", "–40 to –60%", ACCENT_GREEN),
        ("Time-to-Insight", "–50 to –70%", ACCENT_TEAL),
        ("Time to Migrate", "–50%", ACCENT_ORANGE),
        ("Payback Period", "9–14 mo.", ACCENT_PURPLE),
    ]
    card_left = Inches(9.0)
    card_top_start = Inches(1.0)
    card_w = Inches(3.8)
    card_h = Inches(1.05)
    card_gap = Inches(0.15)

    for i, (title, value, color) in enumerate(kpi_data):
        _add_kpi_card(slide, card_left,
                      card_top_start + i * (card_h + card_gap),
                      card_w, card_h, title, value, color)

    # Bottom GitHub Copilot bar
    bar = _add_shape(slide, Inches(0), Inches(6.6), SLIDE_WIDTH, Inches(0.9),
                     DARK_BLUE, MSO_SHAPE.RECTANGLE)
    _add_textbox(slide, Inches(0.5), Inches(6.65), Inches(12), Inches(0.35),
                 "🚀 GitHub Copilot Accelerated",
                 font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.LEFT)

    tags = [
        ("T-SQL → Spark SQL", MICROSOFT_BLUE),
        ("PL/SQL → PySpark", ACCENT_GREEN),
        ("OAC RPD → Semantic Models", ACCENT_ORANGE),
        ("Data Quality", ACCENT_RED),
        ("Refactor & Optimize", ACCENT_TEAL),
        ("Auto Tests", MEDIUM_GRAY),
        ("Doc Generation", MEDIUM_GRAY),
    ]
    tag_left = Inches(0.5)
    for label, color in tags:
        w = Inches(len(label) * 0.11 + 0.3)
        tag = _add_shape(slide, tag_left, Inches(7.05), w, Inches(0.3),
                         color, MSO_SHAPE.ROUNDED_RECTANGLE)
        tag.text_frame.paragraphs[0].text = label
        tag.text_frame.paragraphs[0].font.size = Pt(9)
        tag.text_frame.paragraphs[0].font.color.rgb = WHITE
        tag.text_frame.paragraphs[0].font.bold = True
        tag.text_frame.paragraphs[0].font.name = "Segoe UI"
        tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        tag_left += w + Inches(0.1)


def slide_what_can_be_migrated(prs: Presentation):
    """Slide 3: What Can Be Migrated — 97% coverage breakdown."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "What Can Be Migrated — 97% OAC Object Coverage",
                 font_size=28, bold=True, color=DARK_BLUE)

    _add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
                 "60 of 62 OAC object types are fully automated  •  2 require manual review (cell-level security, custom OAC plugins)",
                 font_size=13, color=MEDIUM_GRAY)

    # ── Left column: automated categories ──
    categories = [
        ("Schema & Data", [
            "Tables, Views, Columns, Data Types → Fabric Lakehouse (Delta)",
            "Materialized Views → Fabric Warehouse MVs",
            "Oracle Mirroring → Fabric Mirroring (near-real-time)",
        ], MICROSOFT_BLUE),
        ("ETL / Data Pipelines", [
            "OAC Data Flows → Dataflow Gen2 + Fabric Notebooks (PySpark)",
            "PL/SQL Stored Procedures → PySpark notebooks",
            "Pivot/Unpivot, Error Row Routing, Parallel Chains",
        ], ACCENT_GREEN),
        ("Semantic Models", [
            "RPD Logical Model → TMDL (tables, relationships, hierarchies)",
            "OAC Expressions → DAX measures (120+ rules)",
            "Calculation Groups, DAX UDFs, Direct Lake on OneLake",
        ], ACCENT_ORANGE),
        ("Reports & Dashboards", [
            "OAC Analyses → PBIR reports (80+ visual types)",
            "BI Publisher → Paginated Reports (.rdl)",
            "Alerts/Agents → Data Activator triggers",
            "Action Links → Translytical Task Flows + Drillthrough",
        ], ACCENT_TEAL),
        ("Security & Governance", [
            "Application Roles → Fabric Workspace Roles",
            "Row-Level Security → DAX RLS (dynamic, hierarchy-based)",
            "Object-Level Security → OLS definitions",
            "AAD Group Provisioning, Audit Trail Migration",
        ], ACCENT_PURPLE),
    ]

    y = Inches(1.3)
    for cat_name, items, color in categories:
        # Category header
        bar = _add_shape(slide, Inches(0.5), y, Inches(6), Inches(0.35),
                         color, MSO_SHAPE.ROUNDED_RECTANGLE)
        bar.text_frame.paragraphs[0].text = f"  {cat_name}"
        bar.text_frame.paragraphs[0].font.size = Pt(13)
        bar.text_frame.paragraphs[0].font.bold = True
        bar.text_frame.paragraphs[0].font.color.rgb = WHITE
        bar.text_frame.paragraphs[0].font.name = "Segoe UI"
        y += Inches(0.4)

        for item in items:
            _add_textbox(slide, Inches(0.7), y, Inches(5.8), Inches(0.25),
                         f"✓  {item}", font_size=10, color=BLACK)
            y += Inches(0.22)
        y += Inches(0.08)

    # ── Right column: coverage chart (using shapes to simulate) ──
    chart_left = Inches(7.2)

    _add_textbox(slide, chart_left, Inches(1.3), Inches(5.5), Inches(0.4),
                 "Coverage by Migration Layer",
                 font_size=16, bold=True, color=DARK_BLUE)

    layers = [
        ("Schema & Data", 100, MICROSOFT_BLUE),
        ("ETL Pipelines", 95, ACCENT_GREEN),
        ("Semantic Models", 98, ACCENT_ORANGE),
        ("Reports & Dashboards", 97, ACCENT_TEAL),
        ("Security & Governance", 95, ACCENT_PURPLE),
        ("Overall", 97, ACCENT_GREEN),
    ]
    bar_y = Inches(1.85)
    bar_max_w = Inches(4.0)

    for label, pct, color in layers:
        _add_textbox(slide, chart_left, bar_y, Inches(2.0), Inches(0.3),
                     label, font_size=11, bold=True, color=BLACK)
        # Background bar
        _add_shape(slide, chart_left + Inches(2.1), bar_y + Inches(0.02),
                   bar_max_w, Inches(0.25), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
        # Fill bar
        fill_w = int(bar_max_w * pct / 100)
        _add_shape(slide, chart_left + Inches(2.1), bar_y + Inches(0.02),
                   fill_w, Inches(0.25), color, MSO_SHAPE.ROUNDED_RECTANGLE)
        # Percentage label
        _add_textbox(slide, chart_left + Inches(2.1) + fill_w + Inches(0.1),
                     bar_y, Inches(0.6), Inches(0.3),
                     f"{pct}%", font_size=12, bold=True, color=color)
        bar_y += Inches(0.38)

    # ── Multi-source connectors box ──
    box_y = Inches(4.5)
    _add_textbox(slide, chart_left, box_y, Inches(5.5), Inches(0.35),
                 "Multi-Source Platform Connectors",
                 font_size=14, bold=True, color=DARK_BLUE)

    sources = [
        ("Oracle Analytics Cloud", "Full (primary)", ACCENT_GREEN),
        ("Tableau", "Full (TWB/TWBX + REST)", ACCENT_GREEN),
        ("IBM Cognos", "Full (report specs + expressions)", ACCENT_GREEN),
        ("Qlik Sense", "Full (load scripts + set analysis)", ACCENT_GREEN),
        ("Essbase", "Full (outlines + calc scripts + MDX)", ACCENT_GREEN),
        ("OBIEE", "Metadata + Catalog API", ACCENT_ORANGE),
    ]
    sy = box_y + Inches(0.4)
    for src, status, color in sources:
        _add_textbox(slide, chart_left + Inches(0.1), sy, Inches(2.2), Inches(0.25),
                     src, font_size=10, bold=True, color=BLACK)
        _add_textbox(slide, chart_left + Inches(2.4), sy, Inches(3.0), Inches(0.25),
                     status, font_size=10, color=color)
        sy += Inches(0.22)


def slide_mapping_matrix(prs: Presentation):
    """Slide 4: OAC & Essbase → Fabric Object Mapping Matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "OAC & Essbase → Fabric Object Mapping Matrix",
                 font_size=28, bold=True, color=DARK_BLUE)

    _add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
                 "185 object mappings across 13 categories  •  57% fully automated  •  31% partial  •  8% manual  •  4% not supported",
                 font_size=13, color=MEDIUM_GRAY)

    # ── Summary KPI cards ──
    kpi_data = [
        ("Automated", "106", ACCENT_GREEN),
        ("Partial", "56", ACCENT_ORANGE),
        ("Manual", "16", RGBColor(0xD1, 0x34, 0x38)),
        ("Not Supported", "7", MEDIUM_GRAY),
    ]
    kx = Inches(0.5)
    for title, val, color in kpi_data:
        _add_kpi_card(slide, kx, Inches(1.2), Inches(1.6), Inches(0.9),
                      title, val, color)
        kx += Inches(1.8)

    # ── Category breakdown table ──
    header = ["Category", "Count", "Automated", "Partial", "Manual", "N/A"]
    rows = [
        header,
        ["Asset Discovery & Inventory", "13", "10", "3", "—", "—"],
        ["Schema & Data Type Mapping", "21", "17", "4", "—", "—"],
        ["RPD Logical Model → Semantic Model", "12", "12", "—", "—", "—"],
        ["Expression Translation (OAC → DAX)", "34", "28", "6", "—", "—"],
        ["Essbase Calc Script → DAX", "19", "9", "2", "8", "—"],
        ["Essbase MDX → DAX", "12", "9", "3", "—", "—"],
        ["Essbase Outline → Semantic Model", "14", "9", "3", "1", "1"],
        ["Visual & Report Mapping", "21", "17", "4", "—", "—"],
        ["Prompts & Slicers", "10", "8", "2", "—", "—"],
        ["Security & Governance", "11", "5", "4", "—", "2"],
        ["ETL & Data Pipeline", "18", "12", "4", "2", "—"],
        ["Smart View → Excel", "7", "—", "5", "—", "2"],
        ["Infrastructure & Deployment", "5", "4", "1", "—", "—"],
    ]
    col_widths = [Inches(3.5), Inches(0.7), Inches(1.1), Inches(0.9), Inches(0.9), Inches(0.7)]
    tbl_shape = _add_table(slide, Inches(0.5), Inches(2.3), Inches(7.8),
                           rows, col_widths, row_height=0.27, font_size=10)
    # Colour the status cells
    table = tbl_shape.table
    status_colors = {2: ACCENT_GREEN, 3: ACCENT_ORANGE, 4: ACCENT_RED, 5: MEDIUM_GRAY}
    for ri in range(1, len(rows)):
        for ci, color in status_colors.items():
            cell = table.cell(ri, ci)
            if cell.text_frame.paragraphs[0].text != "—":
                cell.text_frame.paragraphs[0].font.color.rgb = color
                cell.text_frame.paragraphs[0].font.bold = True

    # ── Right side: coverage by agent ──
    _add_textbox(slide, Inches(8.8), Inches(2.3), Inches(4.2), Inches(0.4),
                 "Automation by Agent",
                 font_size=16, bold=True, color=DARK_BLUE)

    agents = [
        ("01 Discovery", 100, MICROSOFT_BLUE),
        ("02 Schema", 90, ACCENT_GREEN),
        ("03 ETL", 62, ACCENT_ORANGE),
        ("04 Semantic Model", 78, ACCENT_TEAL),
        ("05 Report", 85, ACCENT_PURPLE),
        ("06 Security", 50, ACCENT_RED),
        ("07 Validation", 95, ACCENT_GREEN),
    ]
    bar_y = Inches(2.85)
    bar_max_w = Inches(2.8)
    for label, pct, color in agents:
        _add_textbox(slide, Inches(8.8), bar_y, Inches(1.8), Inches(0.3),
                     label, font_size=10, bold=True, color=BLACK)
        _add_shape(slide, Inches(10.6), bar_y + Inches(0.02),
                   bar_max_w, Inches(0.22), LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
        fill_w = int(bar_max_w * pct / 100)
        _add_shape(slide, Inches(10.6), bar_y + Inches(0.02),
                   fill_w, Inches(0.22), color, MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_textbox(slide, Inches(10.6) + fill_w + Inches(0.05),
                     bar_y, Inches(0.5), Inches(0.25),
                     f"{pct}%", font_size=10, bold=True, color=color)
        bar_y += Inches(0.35)

    # ── Bottom note ──
    _add_shape(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7),
               LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.5),
                 "📄 Full interactive matrix: docs/MIGRATION_MATRIX.html  •  "
                 "Detailed rules: docs/MAPPING_REFERENCE.md (§1–§9)  •  "
                 "Smart View guide: SMART_VIEW_TO_EXCEL_MIGRATION.md",
                 font_size=11, color=DARK_BLUE)


def slide_essbase_smart_view(prs: Presentation):
    """Slide 5: Essbase & Smart View Migration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "Essbase & Smart View → Microsoft Fabric",
                 font_size=28, bold=True, color=DARK_BLUE)

    _add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
                 "End-to-end cube migration: outlines, calc scripts, MDX, Smart View workbooks → Fabric Semantic Models + Excel",
                 font_size=13, color=MEDIUM_GRAY)

    # ── Left: Essbase Outline → TMDL ──
    _add_textbox(slide, Inches(0.5), Inches(1.3), Inches(4), Inches(0.35),
                 "Essbase Outline → Semantic Model",
                 font_size=14, bold=True, color=ACCENT_TEAL)

    outline_rows = [
        ["Essbase Concept", "Fabric/PBI Equivalent"],
        ["Cube", "TMDL Semantic Model"],
        ["Dense Dimension", "Fact table columns"],
        ["Sparse Dimension", "Dimension table + relationship"],
        ["Accounts Dimension", "DAX Measures (unpivoted)"],
        ["Time Dimension", "Calendar table + date hierarchy"],
        ["Generation/Level", "Hierarchy levels"],
        ["Dynamic Calc Member", "DAX Measure"],
        ["Stored Member", "Column / Row"],
        ["Attribute Dimension", "Column on dimension table"],
        ["Substitution Variable", "M Parameter / What-If"],
        ["Essbase Filter (Security)", "RLS Role (DAX filter)"],
    ]
    col_w_o = [Inches(2.0), Inches(2.5)]
    _add_table(slide, Inches(0.3), Inches(1.7), Inches(4.5),
               outline_rows, col_w_o, header_color=ACCENT_TEAL,
               row_height=0.3, font_size=10)

    # ── Center: Calc Script + MDX → DAX ──
    _add_textbox(slide, Inches(5.0), Inches(1.3), Inches(3.5), Inches(0.35),
                 "Calc Scripts & MDX → DAX",
                 font_size=14, bold=True, color=ACCENT_ORANGE)

    calc_rows = [
        ["Essbase", "DAX Equivalent"],
        ["@SUM / @AVG / @COUNT", "SUM / AVERAGE / COUNT"],
        ["@PRIOR / @NEXT", "DATEADD / PREVIOUSMONTH"],
        ["@PARENTVAL / @ANCEST", "CALCULATE + ALL/ALLEXCEPT"],
        ["@ISMBR / @ISLEV", "SELECTEDVALUE / level check"],
        ["IF/ELSEIF/ENDIF", "IF / SWITCH(TRUE())"],
        ["FIX / ENDFIX", "CALCULATE + filter context"],
        ["AGG / CALC ALL", "Implicit (DAX auto-agg)"],
        ["[Dim].CurrentMember", "SELECTEDVALUE(col)"],
        ["YTD / QTD / MTD", "TOTALYTD / TOTALQTD / TOTALMTD"],
        ["PeriodsToDate", "DATESYTD / DATESQTD"],
        ["ParallelPeriod", "SAMEPERIODLASTYEAR"],
    ]
    col_w_c = [Inches(2.0), Inches(2.2)]
    _add_table(slide, Inches(4.8), Inches(1.7), Inches(4.2),
               calc_rows, col_w_c, header_color=ACCENT_ORANGE,
               row_height=0.3, font_size=10)

    # ── Right: Smart View → Excel ──
    _add_textbox(slide, Inches(9.3), Inches(1.3), Inches(3.8), Inches(0.35),
                 "Smart View → Excel + Semantic Model",
                 font_size=14, bold=True, color=ACCENT_PURPLE)

    sv_rows = [
        ["Smart View", "Excel Equivalent"],
        ["HsGetValue()", "CUBEVALUE()"],
        ["HsSetValue()", "Translytical / Power Apps"],
        ["HsAlias()", "CUBEMEMBERPROPERTY()"],
        ["Member Selection", "CUBESET() + CUBEMEMBER()"],
        ["Ad-hoc Grid", "PivotTable (Analyze in Excel)"],
        ["POV Bar", "PivotTable Filters / Slicers"],
        ["Zoom In/Out", "PivotTable expand/collapse"],
        ["Data Form", "Power Apps + Dataverse"],
        ["Suppress #Missing", "PivotTable value filters"],
        ["Subst. Variable", "Named range / What-If param"],
        ["Cascade POV", "Connected slicers"],
    ]
    col_w_sv = [Inches(1.6), Inches(2.2)]
    _add_table(slide, Inches(9.1), Inches(1.7), Inches(3.8),
               sv_rows, col_w_sv, header_color=ACCENT_PURPLE,
               row_height=0.3, font_size=10)

    # ── Bottom: connection methods & write-back ──
    _add_shape(slide, Inches(0), Inches(5.7), SLIDE_WIDTH, Inches(0.05),
               MICROSOFT_BLUE, MSO_SHAPE.RECTANGLE)

    _add_textbox(slide, Inches(0.5), Inches(5.85), Inches(4), Inches(0.35),
                 "Excel Connection Methods",
                 font_size=13, bold=True, color=DARK_BLUE)
    methods = [
        "✓  Analyze in Excel (Power BI Service — recommended)",
        "✓  XMLA endpoint (power users / automation — F2+)",
        "✓  Power BI Excel add-in (Office 365 — embedded visuals)",
    ]
    my = Inches(6.15)
    for m in methods:
        _add_textbox(slide, Inches(0.7), my, Inches(4.5), Inches(0.22),
                     m, font_size=10, color=BLACK)
        my += Inches(0.22)

    _add_textbox(slide, Inches(5.5), Inches(5.85), Inches(4), Inches(0.35),
                 "Write-Back Options",
                 font_size=13, bold=True, color=DARK_BLUE)
    wb_opts = [
        "✓  Translytical Task Flow (UDF → SQL — Preview)",
        "✓  Power Apps visual (embedded form in report)",
        "✓  Excel + Power Automate (HTTP trigger → SQL)",
    ]
    wy = Inches(6.15)
    for w in wb_opts:
        _add_textbox(slide, Inches(5.7), wy, Inches(4.5), Inches(0.22),
                     w, font_size=10, color=BLACK)
        wy += Inches(0.22)

    _add_textbox(slide, Inches(10.5), Inches(5.85), Inches(2.5), Inches(0.35),
                 "Coverage",
                 font_size=13, bold=True, color=DARK_BLUE)
    cov = [
        ("Read (CUBEVALUE)", "✅ Full", ACCENT_GREEN),
        ("Write-back", "🟡 Partial", ACCENT_ORANGE),
        ("Data Forms", "🟡 Partial", ACCENT_ORANGE),
    ]
    cy = Inches(6.15)
    for label, status, color in cov:
        _add_textbox(slide, Inches(10.5), cy, Inches(1.3), Inches(0.22),
                     label, font_size=10, color=BLACK)
        _add_textbox(slide, Inches(11.8), cy, Inches(1.0), Inches(0.22),
                     status, font_size=10, bold=True, color=color)
        cy += Inches(0.22)


def slide_longview_writeback(prs: Presentation):
    """Slide 6: Longview/EPM — Two migration options comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "Longview / EPM Writeback → Two Migration Paths",
                 font_size=28, bold=True, color=DARK_BLUE)

    _add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.35),
                 "Choose the path that fits your timeline and budget planning maturity",
                 font_size=13, color=MEDIUM_GRAY)

    mid = Inches(6.55)  # vertical divider X

    # ═══════════════════════════════════════════════════════════════════
    # OPTION A — Keep Longview, replace backend
    # ═══════════════════════════════════════════════════════════════════
    opt_a_left = Inches(0.3)
    opt_a_w = Inches(6.0)

    # Header card
    hdr_a = _add_shape(slide, opt_a_left, Inches(1.2), opt_a_w, Inches(0.5),
                       ACCENT_TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
    hdr_a.text_frame.paragraphs[0].text = "Option A — Keep Longview, Replace Backend"
    hdr_a.text_frame.paragraphs[0].font.size = Pt(14)
    hdr_a.text_frame.paragraphs[0].font.bold = True
    hdr_a.text_frame.paragraphs[0].font.color.rgb = WHITE
    hdr_a.text_frame.paragraphs[0].font.name = "Segoe UI"
    hdr_a.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    opt_a_rows = [
        ["Layer", "Before", "After"],
        ["Frontend", "Longview UI", "Longview UI (unchanged)"],
        ["Writeback", "Essbase cube", "Fabric Warehouse (TDS)"],
        ["Calc engine", "Essbase scripts", "Fabric Notebook (PySpark)"],
        ["Reporting", "OAC / Smart View", "Power BI DirectLake"],
        ["Security", "Essbase filters", "Entra ID + RLS"],
        ["License", "Oracle Essbase", "Fabric F SKU"],
    ]
    col_w_a = [Inches(1.2), Inches(1.8), Inches(2.5)]
    _add_table(slide, opt_a_left, Inches(1.85), Inches(5.5),
               opt_a_rows, col_w_a, header_color=ACCENT_TEAL,
               row_height=0.28, font_size=10)

    # Flow diagram
    flow_a = [
        ("Longview", ACCENT_PURPLE),
        ("Fabric\nWarehouse", MICROSOFT_BLUE),
        ("PySpark\nNotebook", ACCENT_ORANGE),
        ("DirectLake\nPower BI", ACCENT_GREEN),
    ]
    bx = opt_a_left + Inches(0.1)
    bw = Inches(1.25)
    for label, color in flow_a:
        card = _add_shape(slide, bx, Inches(4.15), bw, Inches(0.7), color,
                          MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = card.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = "Segoe UI"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bx += bw + Inches(0.15)

    # Arrows
    for i in range(3):
        ax = opt_a_left + Inches(0.1) + (bw + Inches(0.15)) * i + bw
        _add_textbox(slide, ax, Inches(4.25), Inches(0.15), Inches(0.4),
                     "→", font_size=14, bold=True, color=DARK_BLUE,
                     alignment=PP_ALIGN.CENTER)

    # Pros / Cons
    _add_textbox(slide, opt_a_left, Inches(5.0), opt_a_w, Inches(0.25),
                 "✅ Pros", font_size=11, bold=True, color=ACCENT_GREEN)
    pros_a = [
        "Zero change for Longview users (same UI/forms)",
        "Preserves all Longview workflows & approvals",
        "Connection string change only — low risk",
        "Faster rollout (2–4 weeks)",
    ]
    py = Inches(5.25)
    for p in pros_a:
        _add_textbox(slide, opt_a_left + Inches(0.15), py, Inches(5.8), Inches(0.2),
                     f"• {p}", font_size=9, color=BLACK)
        py += Inches(0.18)

    _add_textbox(slide, opt_a_left, py + Inches(0.05), opt_a_w, Inches(0.25),
                 "⚠️ Cons", font_size=11, bold=True, color=ACCENT_ORANGE)
    cons_a = [
        "Keeps Longview license cost",
        "Two tools to maintain (Longview + Power BI)",
    ]
    cy = py + Inches(0.28)
    for c in cons_a:
        _add_textbox(slide, opt_a_left + Inches(0.15), cy, Inches(5.8), Inches(0.2),
                     f"• {c}", font_size=9, color=BLACK)
        cy += Inches(0.18)

    # ═══════════════════════════════════════════════════════════════════
    # Vertical divider
    # ═══════════════════════════════════════════════════════════════════
    _add_shape(slide, mid, Inches(1.2), Inches(0.04), Inches(5.6),
               MICROSOFT_BLUE, MSO_SHAPE.RECTANGLE)
    # "VS" badge
    vs_badge = _add_shape(slide, mid - Inches(0.25), Inches(3.75),
                          Inches(0.55), Inches(0.55), MICROSOFT_BLUE,
                          MSO_SHAPE.OVAL)
    vs_badge.text_frame.paragraphs[0].text = "VS"
    vs_badge.text_frame.paragraphs[0].font.size = Pt(12)
    vs_badge.text_frame.paragraphs[0].font.bold = True
    vs_badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    vs_badge.text_frame.paragraphs[0].font.name = "Segoe UI"
    vs_badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ═══════════════════════════════════════════════════════════════════
    # OPTION B — Drop Longview, use Power BI Writeback on SQL
    # ═══════════════════════════════════════════════════════════════════
    opt_b_left = Inches(6.9)
    opt_b_w = Inches(6.0)

    hdr_b = _add_shape(slide, opt_b_left, Inches(1.2), opt_b_w, Inches(0.5),
                       ACCENT_ORANGE, MSO_SHAPE.ROUNDED_RECTANGLE)
    hdr_b.text_frame.paragraphs[0].text = "Option B — Power BI Writeback on SQL"
    hdr_b.text_frame.paragraphs[0].font.size = Pt(14)
    hdr_b.text_frame.paragraphs[0].font.bold = True
    hdr_b.text_frame.paragraphs[0].font.color.rgb = WHITE
    hdr_b.text_frame.paragraphs[0].font.name = "Segoe UI"
    hdr_b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    opt_b_rows = [
        ["Layer", "Before", "After"],
        ["Frontend", "Longview UI", "Power BI + Power Apps"],
        ["Writeback", "Essbase cube", "Azure SQL / Fabric WH"],
        ["Calc engine", "Essbase scripts", "SQL SP + DAX measures"],
        ["Reporting", "OAC / Smart View", "Power BI DirectLake"],
        ["Security", "Essbase filters", "Entra ID + RLS"],
        ["License", "Oracle + Longview", "Fabric + M365 (no Longview)"],
    ]
    col_w_b = [Inches(1.2), Inches(1.8), Inches(2.5)]
    _add_table(slide, opt_b_left, Inches(1.85), Inches(5.5),
               opt_b_rows, col_w_b, header_color=ACCENT_ORANGE,
               row_height=0.28, font_size=10)

    # Flow diagram
    flow_b = [
        ("Power BI\n+ Power Apps", ACCENT_PURPLE),
        ("Azure SQL /\nFabric WH", MICROSOFT_BLUE),
        ("SQL SPs\n+ DAX", ACCENT_ORANGE),
        ("DirectLake\nPower BI", ACCENT_GREEN),
    ]
    bx2 = opt_b_left + Inches(0.1)
    for label, color in flow_b:
        card = _add_shape(slide, bx2, Inches(4.15), bw, Inches(0.7), color,
                          MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = card.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.name = "Segoe UI"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bx2 += bw + Inches(0.15)

    # Arrows
    for i in range(3):
        ax = opt_b_left + Inches(0.1) + (bw + Inches(0.15)) * i + bw
        _add_textbox(slide, ax, Inches(4.25), Inches(0.15), Inches(0.4),
                     "→", font_size=14, bold=True, color=DARK_BLUE,
                     alignment=PP_ALIGN.CENTER)

    # Pros / Cons
    _add_textbox(slide, opt_b_left, Inches(5.0), opt_b_w, Inches(0.25),
                 "✅ Pros", font_size=11, bold=True, color=ACCENT_GREEN)
    pros_b = [
        "Eliminates both Oracle AND Longview licensing",
        "Single platform (Microsoft stack end-to-end)",
        "Power Apps forms for budget input (modern UX)",
        "Lower total cost of ownership long-term",
    ]
    py2 = Inches(5.25)
    for p in pros_b:
        _add_textbox(slide, opt_b_left + Inches(0.15), py2, Inches(5.8), Inches(0.2),
                     f"• {p}", font_size=9, color=BLACK)
        py2 += Inches(0.18)

    _add_textbox(slide, opt_b_left, py2 + Inches(0.05), opt_b_w, Inches(0.25),
                 "⚠️ Cons", font_size=11, bold=True, color=ACCENT_ORANGE)
    cons_b = [
        "Longview forms/workflows must be rebuilt in Power Apps",
        "Longer migration (6–10 weeks)",
        "User retraining required",
    ]
    cy2 = py2 + Inches(0.28)
    for c in cons_b:
        _add_textbox(slide, opt_b_left + Inches(0.15), cy2, Inches(5.8), Inches(0.2),
                     f"• {c}", font_size=9, color=BLACK)
        cy2 += Inches(0.18)

    # ═══════════════════════════════════════════════════════════════════
    # Bottom recommendation bar
    # ═══════════════════════════════════════════════════════════════════
    rec_bar = _add_shape(slide, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.45),
                         DARK_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    rec_bar.text_frame.paragraphs[0].text = (
        "💡 Recommendation:  Option A for quick wins (keep Longview, swap Essbase in weeks)  "
        "→  Option B as a phased follow-up (retire Longview, consolidate on Power BI)"
    )
    rec_bar.text_frame.paragraphs[0].font.size = Pt(11)
    rec_bar.text_frame.paragraphs[0].font.color.rgb = WHITE
    rec_bar.text_frame.paragraphs[0].font.name = "Segoe UI"
    rec_bar.text_frame.paragraphs[0].font.bold = True
    rec_bar.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def slide_longview_migration_plan(prs: Presentation):
    """Slide 7: Longview — Migration Plans & Time Savings for both options."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "Longview Migration Plans — Timeline & Savings",
                 font_size=28, bold=True, color=DARK_BLUE)

    _add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.35),
                 "Phase-by-phase rollout for each option with effort and savings estimates",
                 font_size=13, color=MEDIUM_GRAY)

    mid = Inches(6.55)

    # ═══════════════════════════════════════════════════════════════════
    # OPTION A — Migration Plan
    # ═══════════════════════════════════════════════════════════════════
    opt_a_left = Inches(0.3)

    hdr_a = _add_shape(slide, opt_a_left, Inches(1.15), Inches(6.0), Inches(0.45),
                       ACCENT_TEAL, MSO_SHAPE.ROUNDED_RECTANGLE)
    hdr_a.text_frame.paragraphs[0].text = "Option A — Keep Longview  (2–4 weeks)"
    hdr_a.text_frame.paragraphs[0].font.size = Pt(13)
    hdr_a.text_frame.paragraphs[0].font.bold = True
    hdr_a.text_frame.paragraphs[0].font.color.rgb = WHITE
    hdr_a.text_frame.paragraphs[0].font.name = "Segoe UI"
    hdr_a.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    plan_a_rows = [
        ["Phase", "Task", "Duration"],
        ["1", "Assessment + Essbase outline inventory", "2 days"],
        ["2", "Fabric Warehouse DDL (auto-generated)", "1 day"],
        ["3", "Data migration (Essbase → Delta tables)", "3–5 days"],
        ["4", "PySpark notebooks (calc script conversion)", "3–5 days"],
        ["5", "Connection string swap + Longview config", "1 day"],
        ["6", "UAT + parallel run with Essbase", "3–5 days"],
        ["7", "Go-live + monitoring", "2 days"],
    ]
    col_w_a = [Inches(0.6), Inches(3.6), Inches(1.2)]
    _add_table(slide, opt_a_left, Inches(1.7), Inches(5.5),
               plan_a_rows, col_w_a, header_color=ACCENT_TEAL,
               row_height=0.27, font_size=9)

    # Gantt-style visual (simplified horizontal bars)
    _add_textbox(slide, opt_a_left, Inches(4.15), Inches(5.5), Inches(0.25),
                 "Timeline", font_size=10, bold=True, color=DARK_BLUE)
    weeks_a = [
        ("Wk 1: Assess + DDL + Data", ACCENT_TEAL, 2.5),
        ("Wk 2: PySpark + Config", MICROSOFT_BLUE, 2.5),
        ("Wk 3: UAT", ACCENT_ORANGE, 1.8),
        ("Wk 4: Go-live", ACCENT_GREEN, 1.2),
    ]
    gy = Inches(4.4)
    for label, color, w in weeks_a:
        bar = _add_shape(slide, opt_a_left + Inches(0.05), gy, Inches(w), Inches(0.28),
                         color, MSO_SHAPE.ROUNDED_RECTANGLE)
        bar.text_frame.paragraphs[0].text = label
        bar.text_frame.paragraphs[0].font.size = Pt(8)
        bar.text_frame.paragraphs[0].font.color.rgb = WHITE
        bar.text_frame.paragraphs[0].font.name = "Segoe UI"
        bar.text_frame.paragraphs[0].font.bold = True
        bar.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        gy += Inches(0.32)

    # ═══════════════════════════════════════════════════════════════════
    # OPTION B — Migration Plan
    # ═══════════════════════════════════════════════════════════════════
    opt_b_left = Inches(6.9)

    hdr_b = _add_shape(slide, opt_b_left, Inches(1.15), Inches(6.0), Inches(0.45),
                       ACCENT_ORANGE, MSO_SHAPE.ROUNDED_RECTANGLE)
    hdr_b.text_frame.paragraphs[0].text = "Option B — Power BI Writeback  (6–10 weeks)"
    hdr_b.text_frame.paragraphs[0].font.size = Pt(13)
    hdr_b.text_frame.paragraphs[0].font.bold = True
    hdr_b.text_frame.paragraphs[0].font.color.rgb = WHITE
    hdr_b.text_frame.paragraphs[0].font.name = "Segoe UI"
    hdr_b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    plan_b_rows = [
        ["Phase", "Task", "Duration"],
        ["1", "Assessment + writeback requirements", "3–5 days"],
        ["2", "Azure SQL / Fabric WH tables + SQL SPs", "5–7 days"],
        ["3", "Power Apps input forms design & build", "5–10 days"],
        ["4", "DAX measures + DirectLake semantic model", "3–5 days"],
        ["5", "Power BI reports + writeback visuals", "5–7 days"],
        ["6", "User training + UAT", "5–7 days"],
        ["7", "Go-live + Longview decommission plan", "3–5 days"],
    ]
    col_w_b = [Inches(0.6), Inches(3.6), Inches(1.2)]
    _add_table(slide, opt_b_left, Inches(1.7), Inches(5.5),
               plan_b_rows, col_w_b, header_color=ACCENT_ORANGE,
               row_height=0.27, font_size=9)

    # Gantt-style visual
    _add_textbox(slide, opt_b_left, Inches(4.15), Inches(5.5), Inches(0.25),
                 "Timeline", font_size=10, bold=True, color=DARK_BLUE)
    weeks_b = [
        ("Wk 1–2: Assess + SQL tables + SPs", ACCENT_ORANGE, 3.2),
        ("Wk 3–5: Power Apps + DAX + Reports", MICROSOFT_BLUE, 4.0),
        ("Wk 6–8: Training + UAT", ACCENT_PURPLE, 3.2),
        ("Wk 9–10: Go-live + Decommission", ACCENT_GREEN, 2.5),
    ]
    gy2 = Inches(4.4)
    for label, color, w in weeks_b:
        bar = _add_shape(slide, opt_b_left + Inches(0.05), gy2, Inches(w), Inches(0.28),
                         color, MSO_SHAPE.ROUNDED_RECTANGLE)
        bar.text_frame.paragraphs[0].text = label
        bar.text_frame.paragraphs[0].font.size = Pt(8)
        bar.text_frame.paragraphs[0].font.color.rgb = WHITE
        bar.text_frame.paragraphs[0].font.name = "Segoe UI"
        bar.text_frame.paragraphs[0].font.bold = True
        bar.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        gy2 += Inches(0.32)

    # ═══════════════════════════════════════════════════════════════════
    # Vertical divider
    # ═══════════════════════════════════════════════════════════════════
    _add_shape(slide, mid, Inches(1.15), Inches(0.04), Inches(4.55),
               MICROSOFT_BLUE, MSO_SHAPE.RECTANGLE)

    # ═══════════════════════════════════════════════════════════════════
    # Bottom — Time & Cost Savings comparison table
    # ═══════════════════════════════════════════════════════════════════
    _add_textbox(slide, Inches(0.5), Inches(5.65), Inches(6), Inches(0.3),
                 "Time & Cost Savings vs. Current State",
                 font_size=14, bold=True, color=DARK_BLUE)

    savings_rows = [
        ["Metric", "Current (Oracle + Longview)", "Option A", "Option B"],
        ["Migration Duration", "—", "2–4 weeks", "6–10 weeks"],
        ["User Disruption", "—", "Zero (same Longview UI)", "Moderate (new Power Apps)"],
        ["Oracle License Eliminated", "—", "✅ Yes", "✅ Yes"],
        ["Longview License Eliminated", "—", "❌ No", "✅ Yes"],
        ["Ongoing Maintenance", "6+ FTEs", "2–3 FTEs", "1–2 FTEs"],
        ["Annual Savings (est.)", "Baseline", "$200K–$500K/yr", "$400K–$800K/yr"],
    ]
    col_w_s = [Inches(2.2), Inches(2.8), Inches(2.8), Inches(2.8)]
    tbl = _add_table(slide, Inches(0.5), Inches(5.95), Inches(12.3),
                     savings_rows, col_w_s, header_color=DARK_BLUE,
                     row_height=0.24, font_size=9)
    # Highlight the savings row
    table = tbl.table
    last = len(savings_rows) - 1
    for ci in range(4):
        c = table.cell(last, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT_GREEN
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.text_frame.paragraphs[0].font.bold = True


def slide_time_savings(prs: Presentation):
    """Slide 8: Automation & Time Savings."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "Automation & Time Savings",
                 font_size=28, bold=True, color=DARK_BLUE)

    # ── Before vs After comparison ──
    header = ["Migration Phase", "Manual (Traditional)", "Automated (This Tool)", "Time Saved"]
    rows = [
        header,
        ["Discovery & Inventory", "2–4 weeks", "< 1 day", "90%+"],
        ["Schema Migration (DDL)", "1–2 weeks", "Minutes (auto-generated)", "95%+"],
        ["ETL Pipeline Conversion", "4–8 weeks", "1–2 weeks", "60–75%"],
        ["Semantic Model (RPD → TMDL)", "3–6 weeks", "1–3 days", "85%+"],
        ["Report Migration", "4–8 weeks", "1–2 weeks", "60–75%"],
        ["Security (RLS/OLS)", "1–2 weeks", "Hours", "90%+"],
        ["Validation & Testing", "2–4 weeks", "< 1 week (4,005 auto tests)", "70%+"],
        ["Total End-to-End", "8–12 months", "4–6 months", "~50%"],
    ]
    col_widths = [Inches(2.4), Inches(2.5), Inches(3.2), Inches(1.2)]
    tbl = _add_table(slide, Inches(0.5), Inches(1.0), Inches(9.3), rows, col_widths)
    # Highlight total row
    table = tbl.table
    last = len(rows) - 1
    for ci in range(4):
        c = table.cell(last, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT_GREEN
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.text_frame.paragraphs[0].font.bold = True

    # ── Right side: key stats ──
    stats = [
        ("188+", "Source Modules", MICROSOFT_BLUE),
        ("4,005", "Automated Tests", ACCENT_GREEN),
        ("185", "Object Mappings", ACCENT_ORANGE),
        ("80+", "Visual Type Mappings", ACCENT_TEAL),
        ("8", "Specialized Agents", ACCENT_PURPLE),
    ]
    sx = Inches(10.2)
    sy = Inches(1.0)
    for val, label, color in stats:
        card = _add_shape(slide, sx, sy, Inches(2.8), Inches(1.0), WHITE)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        _add_textbox(slide, sx + Inches(0.1), sy + Inches(0.05),
                     Inches(2.6), Inches(0.55),
                     val, font_size=28, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, sx + Inches(0.1), sy + Inches(0.55),
                     Inches(2.6), Inches(0.35),
                     label, font_size=11, color=MEDIUM_GRAY,
                     alignment=PP_ALIGN.CENTER)
        sy += Inches(1.15)

    # ── Bottom insight ──
    insight = _add_shape(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8),
                         LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_textbox(slide, Inches(0.8), Inches(6.1), Inches(11.8), Inches(0.6),
                 "💡 GitHub Copilot further accelerates code translation tasks (PL/SQL → PySpark, OAC expressions → DAX, "
                 "Essbase calc scripts → DAX) by 30–50%, reducing the \"With Copilot\" timelines shown in the project matrix.",
                 font_size=12, color=DARK_BLUE)


def slide_roi(prs: Presentation):
    """Slide 8: ROI & Business Impact."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "ROI & Business Impact",
                 font_size=28, bold=True, color=DARK_BLUE)

    # ── Large KPI cards across top ──
    top_kpis = [
        ("Infrastructure TCO", "–50 to –65%", "Eliminate Oracle licensing,\nhardware, & DBA overhead", MICROSOFT_BLUE),
        ("Operational Effort", "–40 to –60%", "Automated pipelines replace\nmanual ETL & report builds", ACCENT_GREEN),
        ("Time to Migrate", "–50%", "AI-accelerated migration\n8–12 mo → 4–6 mo", ACCENT_ORANGE),
        ("Time to Insight", "–50 to –70%", "Self-service analytics,\nreal-time dashboards", ACCENT_TEAL),
    ]
    kw = Inches(2.95)
    kh = Inches(1.5)
    kx = Inches(0.5)
    for title, value, detail, color in top_kpis:
        card = _add_shape(slide, kx, Inches(1.0), kw, kh, WHITE)
        card.line.color.rgb = color
        card.line.width = Pt(3)
        _add_textbox(slide, kx + Inches(0.1), Inches(1.05), kw - Inches(0.2), Inches(0.3),
                     title, font_size=12, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, kx + Inches(0.1), Inches(1.35), kw - Inches(0.2), Inches(0.5),
                     value, font_size=26, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, kx + Inches(0.1), Inches(1.9), kw - Inches(0.2), Inches(0.5),
                     detail, font_size=10, color=MEDIUM_GRAY,
                     alignment=PP_ALIGN.CENTER)
        kx += kw + Inches(0.15)

    # ── Cost comparison table ──
    _add_textbox(slide, Inches(0.5), Inches(2.8), Inches(6), Inches(0.4),
                 "3-Year Total Cost of Ownership (indicative)",
                 font_size=16, bold=True, color=DARK_BLUE)

    cost_rows = [
        ["Cost Category", "Oracle (Current)", "Microsoft Fabric", "Savings"],
        ["Database Licensing", "$800K–$1.5M/yr", "Included in Fabric", "100%"],
        ["BI Licensing (OAC)", "$200K–$500K/yr", "Power BI Pro/Premium", "40–60%"],
        ["Infrastructure (HW/Cloud)", "$300K–$600K/yr", "Pay-per-use (CU)", "50–70%"],
        ["DBA / ETL Admin", "4–6 FTEs", "1–2 FTEs", "60–75%"],
        ["Report Development", "3–5 FTEs", "1–2 FTEs (self-service)", "50–70%"],
        ["3-Year Net Savings", "", "", "$2M–$6M+"],
    ]
    col_widths_cost = [Inches(2.0), Inches(2.0), Inches(2.2), Inches(1.2)]
    tbl = _add_table(slide, Inches(0.5), Inches(3.3), Inches(7.4),
                     cost_rows, col_widths_cost)
    table = tbl.table
    last = len(cost_rows) - 1
    for ci in range(4):
        c = table.cell(last, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = ACCENT_GREEN
        c.text_frame.paragraphs[0].font.color.rgb = WHITE
        c.text_frame.paragraphs[0].font.bold = True

    # ── Right: payback timeline ──
    _add_textbox(slide, Inches(8.3), Inches(2.8), Inches(4.5), Inches(0.4),
                 "Payback Timeline",
                 font_size=16, bold=True, color=DARK_BLUE)

    milestones = [
        ("Month 1–2", "Discovery + Schema Migration", MICROSOFT_BLUE),
        ("Month 3–4", "ETL + Semantic Model + Reports", ACCENT_GREEN),
        ("Month 5–6", "Security + Validation + Go-Live", ACCENT_ORANGE),
        ("Month 9–14", "Full ROI Payback", ACCENT_PURPLE),
    ]
    my = Inches(3.3)
    for period, desc, color in milestones:
        # Timeline dot
        dot = _add_shape(slide, Inches(8.5), my + Inches(0.05),
                         Inches(0.2), Inches(0.2), color, MSO_SHAPE.OVAL)
        # Line
        if period != "Month 9–14":
            _add_shape(slide, Inches(8.58), my + Inches(0.25),
                       Inches(0.04), Inches(0.6), color, MSO_SHAPE.RECTANGLE)
        _add_textbox(slide, Inches(8.9), my, Inches(1.5), Inches(0.3),
                     period, font_size=12, bold=True, color=color)
        _add_textbox(slide, Inches(8.9), my + Inches(0.25), Inches(4.0), Inches(0.3),
                     desc, font_size=10, color=MEDIUM_GRAY)
        my += Inches(0.7)

    # ── Bottom: risk mitigation ──
    _add_textbox(slide, Inches(0.5), Inches(6.1), Inches(12), Inches(0.4),
                 "Risk Mitigation",
                 font_size=14, bold=True, color=DARK_BLUE)
    risks = ("✓ Incremental migration (wave-based)    "
             "✓ Automated rollback    "
             "✓ 4,005 validation tests    "
             "✓ Data reconciliation (row counts + checksums)    "
             "✓ Visual regression testing")
    _add_textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
                 risks, font_size=11, color=ACCENT_GREEN)


def slide_platform_architecture(prs: Presentation):
    """Slide 9: Platform Architecture (simplified)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)

    _add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
                 "Migration Platform Architecture",
                 font_size=28, bold=True, color=DARK_BLUE)

    _add_textbox(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.4),
                 "8 specialized AI agents orchestrated in a DAG pipeline  •  188+ Python modules  •  Zero external dependencies for core",
                 font_size=13, color=MEDIUM_GRAY)

    # ── Source platforms ──
    _add_textbox(slide, Inches(0.3), Inches(1.4), Inches(2), Inches(0.35),
                 "SOURCE PLATFORMS", font_size=12, bold=True, color=MEDIUM_GRAY)
    sources = ["Oracle Analytics Cloud", "Oracle Exadata", "Tableau", "IBM Cognos",
               "Qlik Sense", "Essbase"]
    sy = Inches(1.8)
    for src in sources:
        s = _add_shape(slide, Inches(0.3), sy, Inches(2.2), Inches(0.35),
                       MICROSOFT_BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
        s.text_frame.paragraphs[0].text = src
        s.text_frame.paragraphs[0].font.size = Pt(10)
        s.text_frame.paragraphs[0].font.color.rgb = WHITE
        s.text_frame.paragraphs[0].font.name = "Segoe UI"
        s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        sy += Inches(0.4)

    # ── Arrow ──
    _add_textbox(slide, Inches(2.7), Inches(2.8), Inches(0.5), Inches(0.5),
                 "→", font_size=30, bold=True, color=ACCENT_GREEN,
                 alignment=PP_ALIGN.CENTER)

    # ── Agent pipeline (center) ──
    _add_textbox(slide, Inches(3.3), Inches(1.4), Inches(3.5), Inches(0.35),
                 "8-AGENT PIPELINE", font_size=12, bold=True, color=MEDIUM_GRAY)
    agents = [
        ("01 Discovery", MICROSOFT_BLUE),
        ("02 Schema", ACCENT_GREEN),
        ("03 ETL", ACCENT_ORANGE),
        ("04 Semantic Model", ACCENT_TEAL),
        ("05 Report", ACCENT_PURPLE),
        ("06 Security", ACCENT_RED),
        ("07 Validation", RGBColor(0x88, 0x88, 0x00)),
        ("08 Orchestrator", DARK_BLUE),
    ]
    ay = Inches(1.8)
    for name, color in agents:
        a = _add_shape(slide, Inches(3.5), ay, Inches(2.2), Inches(0.35),
                       color, MSO_SHAPE.ROUNDED_RECTANGLE)
        a.text_frame.paragraphs[0].text = name
        a.text_frame.paragraphs[0].font.size = Pt(10)
        a.text_frame.paragraphs[0].font.color.rgb = WHITE
        a.text_frame.paragraphs[0].font.name = "Segoe UI"
        a.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ay += Inches(0.4)

    # ── Arrow ──
    _add_textbox(slide, Inches(5.9), Inches(2.8), Inches(0.5), Inches(0.5),
                 "→", font_size=30, bold=True, color=ACCENT_GREEN,
                 alignment=PP_ALIGN.CENTER)

    # ── AI translation engine (center-right) ──
    _add_textbox(slide, Inches(6.5), Inches(1.4), Inches(3), Inches(0.35),
                 "AI TRANSLATION ENGINE", font_size=12, bold=True, color=MEDIUM_GRAY)
    translations = [
        "Oracle DDL → Delta Tables",
        "PL/SQL → PySpark",
        "OAC Expressions → DAX (120+ rules)",
        "Essbase Calc/MDX → DAX (75+ rules)",
        "RPD Model → TMDL",
        "OAC Visuals → PBIR (80+ types)",
        "Smart View → Excel CUBE functions",
        "Session Vars → DAX RLS",
    ]
    ty = Inches(1.8)
    for t in translations:
        tb = _add_shape(slide, Inches(6.5), ty, Inches(3.3), Inches(0.35),
                        LIGHT_GRAY, MSO_SHAPE.ROUNDED_RECTANGLE)
        tb.text_frame.paragraphs[0].text = t
        tb.text_frame.paragraphs[0].font.size = Pt(9)
        tb.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        tb.text_frame.paragraphs[0].font.name = "Segoe UI"
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ty += Inches(0.4)

    # ── Arrow ──
    _add_textbox(slide, Inches(10.0), Inches(2.8), Inches(0.5), Inches(0.5),
                 "→", font_size=30, bold=True, color=ACCENT_GREEN,
                 alignment=PP_ALIGN.CENTER)

    # ── Target (right) ──
    _add_textbox(slide, Inches(10.5), Inches(1.4), Inches(2.5), Inches(0.35),
                 "FABRIC TARGETS", font_size=12, bold=True, color=MEDIUM_GRAY)
    targets = [
        ("Fabric Lakehouse", ACCENT_GREEN),
        ("Fabric Warehouse", ACCENT_GREEN),
        ("Fabric Pipelines", ACCENT_GREEN),
        ("TMDL Semantic Model", ACCENT_GREEN),
        ("PBIR Reports", ACCENT_GREEN),
        ("Paginated Reports", ACCENT_GREEN),
        ("Data Activator", ACCENT_GREEN),
        ("Power BI Service", ACCENT_GREEN),
    ]
    ty2 = Inches(1.8)
    for tname, color in targets:
        t = _add_shape(slide, Inches(10.5), ty2, Inches(2.2), Inches(0.35),
                       color, MSO_SHAPE.ROUNDED_RECTANGLE)
        t.text_frame.paragraphs[0].text = tname
        t.text_frame.paragraphs[0].font.size = Pt(10)
        t.text_frame.paragraphs[0].font.color.rgb = WHITE
        t.text_frame.paragraphs[0].font.name = "Segoe UI"
        t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        ty2 += Inches(0.4)

    # ── Bottom note ──
    _add_shape(slide, Inches(0), Inches(5.8), SLIDE_WIDTH, Inches(0.05),
               MICROSOFT_BLUE, MSO_SHAPE.RECTANGLE)
    _add_textbox(slide, Inches(0.5), Inches(5.95), Inches(12), Inches(1.2),
                 "Key differentiators:  Hybrid rules-first + LLM translation  •  "
                 "Wave-based incremental migration  •  "
                 "Automated rollback  •  "
                 "Real-time migration dashboard  •  "
                 "4,005 automated validation tests  •  "
                 "185 object mappings  •  "
                 "Direct Lake on OneLake support",
                 font_size=12, color=DARK_BLUE)


def slide_next_steps(prs: Presentation):
    """Slide 10: Next Steps / Call to Action."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BLUE)

    _add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
                 "Next Steps",
                 font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    steps = [
        ("1", "Discovery Workshop", "Run automated inventory against your OAC environment\nto produce a complete asset catalog and complexity assessment.", MICROSOFT_BLUE),
        ("2", "Proof of Concept", "Migrate a representative subset (1 subject area, 5–10 reports)\nto validate fidelity, performance, and ROI assumptions.", ACCENT_GREEN),
        ("3", "Wave Planning", "Use complexity scoring to plan 3–5 migration waves with\nautomated effort estimates and risk-based sequencing.", ACCENT_ORANGE),
        ("4", "Production Migration", "Execute wave-by-wave with automated validation,\nrollback safety nets, and real-time progress dashboards.", ACCENT_TEAL),
    ]

    sx = Inches(0.5)
    card_w = Inches(2.9)
    for num, title, desc, color in steps:
        # Number circle
        circle = _add_shape(slide, sx + Inches(1.1), Inches(2.0),
                            Inches(0.7), Inches(0.7), color, MSO_SHAPE.OVAL)
        circle.text_frame.paragraphs[0].text = num
        circle.text_frame.paragraphs[0].font.size = Pt(24)
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].font.color.rgb = WHITE
        circle.text_frame.paragraphs[0].font.name = "Segoe UI"
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        _add_textbox(slide, sx, Inches(2.9), card_w, Inches(0.4),
                     title, font_size=16, bold=True, color=color,
                     alignment=PP_ALIGN.CENTER)
        _add_textbox(slide, sx, Inches(3.4), card_w, Inches(1.0),
                     desc, font_size=11, color=WHITE,
                     alignment=PP_ALIGN.CENTER)
        sx += card_w + Inches(0.2)

    # Footer
    _add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
                 "Ready to start?   Contact us to schedule a Discovery Workshop.",
                 font_size=20, bold=True, color=ACCENT_TEAL,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
                 "v9.0  •  84 phases complete  •  189+ modules  •  4,045 tests  •  185 object mappings  •  97% OAC coverage  •  MIT License",
                 font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ── Main ───────────────────────────────────────────────────────────────

def build_deck(output_path: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)
    slide_migration_matrix(prs)
    slide_what_can_be_migrated(prs)
    slide_mapping_matrix(prs)
    slide_essbase_smart_view(prs)
    slide_longview_writeback(prs)
    slide_longview_migration_plan(prs)
    slide_time_savings(prs)
    slide_roi(prs)
    slide_platform_architecture(prs)
    slide_next_steps(prs)

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="Generate OAC→Fabric Business Migration ROI deck")
    parser.add_argument("--output", "-o", default="output/OAC_to_Fabric_Business_Deck.pptx",
                        help="Output .pptx path (default: output/OAC_to_Fabric_Business_Deck.pptx)")
    args = parser.parse_args()

    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    path = build_deck(args.output)
    print(f"✅ Business deck generated: {path}")
    print(f"   11 slides • widescreen 16:9 • {os.path.getsize(path):,} bytes")


if __name__ == "__main__":
    main()
