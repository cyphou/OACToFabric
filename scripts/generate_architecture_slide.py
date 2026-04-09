#!/usr/bin/env python3
"""Generate a slide comparing OAC/Essbase architecture vs Fabric/Power BI architecture.

Usage:
    python scripts/generate_architecture_slide.py [--output path/to/slide.pptx]
"""
from __future__ import annotations

import argparse
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x10, 0x7C, 0x10)
TEAL = RGBColor(0x00, 0xB7, 0xC3)
PURPLE = RGBColor(0x6B, 0x2D, 0x8B)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ORACLE_RED = RGBColor(0xC7, 0x40, 0x34)
ORACLE_DARK = RGBColor(0x8B, 0x20, 0x15)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF2, 0xF1)
LIGHT_RED = RGBColor(0xFD, 0xEE, 0xE8)
LIGHT_BLUE = RGBColor(0xE5, 0xF0, 0xFD)
LIGHT_GREEN = RGBColor(0xE5, 0xFD, 0xE8)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ──

def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _shape(slide, l, t, w, h, fill, st=MSO_SHAPE.ROUNDED_RECTANGLE,
           border=None, border_w=Pt(1)):
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.fill.solid()
        s.line.fill.fore_color.rgb = border
        s.line.width = border_w
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _txt(slide, l, t, w, h, text, sz=18, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return tb


def _box(slide, l, t, w, h, fill, label, sub=None, label_sz=10,
         sub_sz=7.5, label_color=WHITE, border=None):
    """Draw a component box with label + optional subtitle."""
    _shape(slide, l, t, w, h, fill, border=border)
    if sub:
        _txt(slide, l, t + Inches(0.03), w, Inches(0.28),
             label, sz=label_sz, bold=True, color=label_color,
             align=PP_ALIGN.CENTER)
        _txt(slide, l, t + Inches(0.26), w, Inches(0.20),
             sub, sz=sub_sz, color=label_color, align=PP_ALIGN.CENTER)
    else:
        _txt(slide, l, t + Inches(0.06), w, Inches(0.32),
             label, sz=label_sz, bold=True, color=label_color,
             align=PP_ALIGN.CENTER)


def _arrow_right(slide, l, t, length, color=GRAY):
    """Horizontal arrow pointing right."""
    conn = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t,
                                  length, Inches(0.22))
    conn.fill.solid()
    conn.fill.fore_color.rgb = color
    conn.line.fill.background()
    conn.shadow.inherit = False


def _arrow_down(slide, l, t, length, color=GRAY):
    """Vertical arrow pointing down."""
    conn = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t,
                                  Inches(0.22), length)
    conn.fill.solid()
    conn.fill.fore_color.rgb = color
    conn.line.fill.background()
    conn.shadow.inherit = False


def _section_label(slide, l, t, w, text, color):
    """Layer label (rotated-style label on the left)."""
    _txt(slide, l, t, w, Inches(0.25), text, sz=7.5, bold=True,
         color=color, align=PP_ALIGN.LEFT)


# ── Diagram builders ──

def _draw_oracle_diagram(slide, x0, y0):
    """Draw the OAC / Essbase / Oracle architecture."""
    W = Inches(5.6)
    H = Inches(5.75)
    # Section background
    _shape(slide, x0, y0, W, H, LIGHT_RED,
           st=MSO_SHAPE.ROUNDED_RECTANGLE, border=ORACLE_RED, border_w=Pt(1.5))

    # Section title
    _shape(slide, x0, y0, W, Inches(0.35), ORACLE_RED, st=MSO_SHAPE.RECTANGLE)
    _txt(slide, x0, y0 + Inches(0.04), W, Inches(0.28),
         "Oracle Analytics Cloud / Essbase", sz=13, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

    # ── Layer 0: OAC Agents & AI ──
    ly0 = y0 + Inches(0.45)
    _section_label(slide, x0 + Inches(0.1), ly0 - Inches(0.02),
                   Inches(1.5), "AGENTS  &  AI", ORACLE_DARK)

    ax = x0 + Inches(0.15)
    aw = Inches(1.2)
    ah = Inches(0.42)
    agap = Inches(0.12)
    _box(slide, ax, ly0 + Inches(0.16), aw, ah, ORACLE_RED,
         "OAC Agents", "Alerts & schedules")
    _box(slide, ax + aw + agap, ly0 + Inches(0.16), aw, ah, ORACLE_RED,
         "Digital Assistant", "NL queries")
    _box(slide, ax + 2 * (aw + agap), ly0 + Inches(0.16), aw, ah, ORACLE_RED,
         "Auto Insights", "ML-driven discovery")
    _box(slide, ax + 3 * (aw + agap), ly0 + Inches(0.16), aw, ah, ORACLE_RED,
         "Ask Oracle", "Natural language")

    # ── Arrow down ──
    _arrow_down(slide, x0 + W / 2 - Inches(0.11), ly0 + Inches(0.62),
                Inches(0.28), ORACLE_RED)

    # ── Layer 1: Users ──
    ly = y0 + Inches(1.25)
    _section_label(slide, x0 + Inches(0.1), ly - Inches(0.02),
                   Inches(1.0), "CONSUMERS", ORACLE_DARK)

    bx = x0 + Inches(0.15)
    bw = Inches(1.2)
    bh = Inches(0.38)
    gap = Inches(0.12)

    _box(slide, bx, ly + Inches(0.18), bw, bh, ORACLE_RED,
         "OAC Analyses", "Interactive reports")
    _box(slide, bx + bw + gap, ly + Inches(0.18), bw, bh, ORACLE_RED,
         "OAC Dashboards", "KPI summaries")
    _box(slide, bx + 2 * (bw + gap), ly + Inches(0.18), bw, bh, ORACLE_RED,
         "Smart View", "Excel add-in")
    _box(slide, bx + 3 * (bw + gap), ly + Inches(0.18), bw, bh, ORACLE_RED,
         "BI Publisher", "Pixel-perfect")

    # ── Arrow down ──
    _arrow_down(slide, x0 + W / 2 - Inches(0.11), ly + Inches(0.62),
                Inches(0.28), ORACLE_RED)

    # ── Layer 2: Semantic ──
    ly2 = ly + Inches(0.98)
    _section_label(slide, x0 + Inches(0.1), ly2 - Inches(0.02),
                   Inches(1.2), "SEMANTIC  LAYER", ORACLE_DARK)

    sw = Inches(2.5)
    sh = Inches(0.42)
    sx = x0 + Inches(0.3)
    _box(slide, sx, ly2 + Inches(0.18), sw, sh, ORACLE_DARK,
         "RPD Physical Layer", "Tables, columns, joins")
    _box(slide, sx + sw + Inches(0.15), ly2 + Inches(0.18), sw, sh, ORACLE_DARK,
         "RPD Logical Layer", "Stars, dims, measures")

    _arrow_down(slide, x0 + W / 2 - Inches(0.11), ly2 + Inches(0.66),
                Inches(0.28), ORACLE_DARK)

    # ── Layer 3: OLAP / ETL ──
    ly3 = ly2 + Inches(1.02)
    _section_label(slide, x0 + Inches(0.1), ly3 - Inches(0.02),
                   Inches(1.2), "OLAP  &  ETL", ORACLE_DARK)

    ow = Inches(1.55)
    oh = Inches(0.42)
    ox = x0 + Inches(0.15)
    ogap = Inches(0.1)
    _box(slide, ox, ly3 + Inches(0.18), ow, oh, ORANGE,
         "Essbase BSO", "Block storage")
    _box(slide, ox + ow + ogap, ly3 + Inches(0.18), ow, oh, ORANGE,
         "Essbase ASO", "Aggregate storage")
    _box(slide, ox + 2 * (ow + ogap), ly3 + Inches(0.18), ow + Inches(0.4), oh,
         ORACLE_DARK, "OAC Data Flows", "ETL pipelines")

    _arrow_down(slide, x0 + W / 2 - Inches(0.11), ly3 + Inches(0.66),
                Inches(0.28), ORACLE_DARK)

    # ── Layer 4: Data ──
    ly4 = ly3 + Inches(1.02)
    _section_label(slide, x0 + Inches(0.1), ly4 - Inches(0.02),
                   Inches(1.2), "DATA  SOURCES", ORACLE_DARK)

    dw = Inches(2.5)
    dh = Inches(0.42)
    dx = x0 + Inches(0.3)
    _box(slide, dx, ly4 + Inches(0.18), dw, dh, RGBColor(0x60, 0x20, 0x10),
         "Oracle Database", "Exadata / RAC")
    _box(slide, dx + dw + Inches(0.15), ly4 + Inches(0.18), dw, dh,
         RGBColor(0x60, 0x20, 0x10),
         "External Sources", "Files, APIs, LDAP")

    # ── Security sidebar ──
    sy = y0 + Inches(0.45)
    _shape(slide, x0 + W - Inches(0.03), sy + Inches(0.15),
           Inches(0.03), Inches(5.0), ORACLE_RED, st=MSO_SHAPE.RECTANGLE)


COPILOT_BLUE = RGBColor(0x6C, 0x2B, 0xD9)  # Copilot purple-blue
LIGHT_PURPLE = RGBColor(0xF0, 0xE8, 0xFD)
O365_RED = RGBColor(0xD8, 0x38, 0x20)

# Logo directory (resolved relative to this script)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_DIR = os.path.join(os.path.dirname(SCRIPT_DIR), "assets", "logos")


def _draw_fabric_diagram(slide, x0, y0):
    """Draw the Microsoft Fabric / Power BI architecture with Data Agents."""
    W = Inches(5.6)
    H = Inches(5.75)
    # Section background
    _shape(slide, x0, y0, W, H, LIGHT_BLUE,
           st=MSO_SHAPE.ROUNDED_RECTANGLE, border=MS_BLUE, border_w=Pt(1.5))

    # Section title
    _shape(slide, x0, y0, W, Inches(0.35), MS_BLUE, st=MSO_SHAPE.RECTANGLE)
    _txt(slide, x0, y0 + Inches(0.04), W, Inches(0.28),
         "Microsoft Fabric & Power BI", sz=13, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

    # ── Layer 0: Copilot & Data Agents + O365 ──
    ly0 = y0 + Inches(0.45)
    _section_label(slide, x0 + Inches(0.1), ly0 - Inches(0.02),
                   Inches(1.5), "COPILOT  &  DATA  AGENTS", COPILOT_BLUE)

    # Agent box
    ax = x0 + Inches(0.15)
    _shape(slide, ax, ly0 + Inches(0.16), Inches(2.1), Inches(0.42),
           COPILOT_BLUE, border=None)
    _txt(slide, ax, ly0 + Inches(0.18), Inches(2.1), Inches(0.22),
         "\U0001F916 Fabric Data Agents", sz=9, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
    _txt(slide, ax, ly0 + Inches(0.38), Inches(2.1), Inches(0.18),
         "Natural language → actions", sz=7, color=WHITE,
         align=PP_ALIGN.CENTER)

    # Arrow from agent to O365 apps
    _arrow_right(slide, ax + Inches(2.2), ly0 + Inches(0.26),
                 Inches(0.35), COPILOT_BLUE)

    # O365 apps with real logos
    app_logos = [
        ("excel", "Excel", RGBColor(0x21, 0x7A, 0x46)),
        ("powerpoint", "PPT", RGBColor(0xD2, 0x4B, 0x27)),
        ("word", "Word", RGBColor(0x29, 0x5A, 0xB8)),
        ("teams", "Teams", RGBColor(0x50, 0x59, 0xC9)),
        ("outlook", "Out", RGBColor(0x00, 0x78, 0xD4)),
    ]
    app_x = ax + Inches(2.65)
    app_w = Inches(0.52)
    app_gap = Inches(0.06)
    for i, (logo_key, name, clr) in enumerate(app_logos):
        lx = app_x + i * (app_w + app_gap)
        logo_path = os.path.join(LOGO_DIR, f"{logo_key}.png")
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                lx + Inches(0.1), ly0 + Inches(0.17),
                Inches(0.32), Inches(0.32))
        else:
            _box(slide, lx, ly0 + Inches(0.16),
                 app_w, Inches(0.42), clr, name, label_sz=7.5)

    # ── Arrow down ──
    _arrow_down(slide, x0 + W / 2 - Inches(0.11), ly0 + Inches(0.62),
                Inches(0.28), COPILOT_BLUE)

    # ── Layer 1: Consumers ──
    ly = y0 + Inches(1.25)
    _section_label(slide, x0 + Inches(0.1), ly - Inches(0.02),
                   Inches(1.0), "CONSUMERS", DARK_BLUE)

    bx = x0 + Inches(0.15)
    bw = Inches(1.2)
    bh = Inches(0.38)
    gap = Inches(0.12)

    _box(slide, bx, ly + Inches(0.18), bw, bh, MS_BLUE,
         "Power BI Reports", "PBIR visuals")
    _box(slide, bx + bw + gap, ly + Inches(0.18), bw, bh, MS_BLUE,
         "PBI Dashboards", "Pinned tiles")
    _box(slide, bx + 2 * (bw + gap), ly + Inches(0.18), bw, bh, MS_BLUE,
         "Analyze in Excel", "CUBE functions")
    _box(slide, bx + 3 * (bw + gap), ly + Inches(0.18), bw, bh, MS_BLUE,
         "Paginated Reports", "Pixel-perfect")

    # ── Arrow down ──
    _arrow_down(slide, x0 + W / 2 - Inches(0.11), ly + Inches(0.62),
                Inches(0.28), MS_BLUE)

    # ── Layer 2: Semantic ──
    ly2 = ly + Inches(0.98)
    _section_label(slide, x0 + Inches(0.1), ly2 - Inches(0.02),
                   Inches(1.2), "SEMANTIC  LAYER", DARK_BLUE)

    sw = Inches(2.5)
    sh = Inches(0.42)
    sx = x0 + Inches(0.3)
    _box(slide, sx, ly2 + Inches(0.18), sw, sh, TEAL,
         "Semantic Model (TMDL)", "Measures, hierarchies")
    _box(slide, sx + sw + Inches(0.15), ly2 + Inches(0.18), sw, sh, TEAL,
         "Direct Lake Mode", "Zero-copy from Delta")

    _arrow_down(slide, x0 + W / 2 - Inches(0.11), ly2 + Inches(0.66),
                Inches(0.28), TEAL)

    # ── Layer 3: Compute / ETL ──
    ly3 = ly2 + Inches(1.02)
    _section_label(slide, x0 + Inches(0.1), ly3 - Inches(0.02),
                   Inches(1.2), "COMPUTE  &  ETL", DARK_BLUE)

    ow = Inches(1.55)
    oh = Inches(0.42)
    ox = x0 + Inches(0.15)
    ogap = Inches(0.1)
    _box(slide, ox, ly3 + Inches(0.18), ow, oh, GREEN,
         "Data Pipelines", "Copy + orchestrate")
    _box(slide, ox + ow + ogap, ly3 + Inches(0.18), ow, oh, GREEN,
         "Dataflow Gen2", "Power Query M")
    _box(slide, ox + 2 * (ow + ogap), ly3 + Inches(0.18), ow + Inches(0.4), oh,
         GREEN, "Spark Notebooks", "PySpark / SQL")

    _arrow_down(slide, x0 + W / 2 - Inches(0.11), ly3 + Inches(0.66),
                Inches(0.28), GREEN)

    # ── Layer 4: Data ──
    ly4 = ly3 + Inches(1.02)
    _section_label(slide, x0 + Inches(0.1), ly4 - Inches(0.02),
                   Inches(1.2), "DATA  STORAGE", DARK_BLUE)

    dw = Inches(2.5)
    dh = Inches(0.42)
    dx = x0 + Inches(0.3)
    _box(slide, dx, ly4 + Inches(0.18), dw, dh, DARK_BLUE,
         "Fabric Lakehouse", "Delta tables (Parquet)")
    _box(slide, dx + dw + Inches(0.15), ly4 + Inches(0.18), dw, dh, DARK_BLUE,
         "Fabric Warehouse", "T-SQL analytics")

    # ── OneLake sidebar ──
    sy = y0 + Inches(0.45)
    _shape(slide, x0 + W - Inches(0.03), sy + Inches(0.15),
           Inches(0.03), Inches(5.0), MS_BLUE, st=MSO_SHAPE.RECTANGLE)


# ── Main ──

def build(output: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)

    # ══ Title bar ══
    _shape(slide, 0, 0, SLIDE_W, Inches(0.9), DARK_BLUE, st=MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.4),
         "Architecture Comparison: Oracle  →  Microsoft Fabric",
         sz=22, bold=True, color=WHITE)
    _txt(slide, Inches(0.5), Inches(0.5), Inches(10), Inches(0.3),
         "From Oracle Analytics Cloud & Essbase to Microsoft Fabric & Power BI",
         sz=11, color=TEAL)

    # ══ Left diagram: Oracle / OAC / Essbase ══
    _draw_oracle_diagram(slide, Inches(0.35), Inches(1.05))

    # ══ Center migration arrow ══
    _arrow_right(slide, Inches(6.15), Inches(4.0), Inches(1.0), TEAL)
    _txt(slide, Inches(6.05), Inches(3.65), Inches(1.2), Inches(0.3),
         "MIGRATE", sz=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ══ Right diagram: Fabric / Power BI ══
    _draw_fabric_diagram(slide, Inches(7.35), Inches(1.05))

    # ══ Legend / footer ══
    fy = Inches(7.0)
    _txt(slide, Inches(0.5), fy, Inches(3), Inches(0.25),
         "■ Consumers   ■ Semantic   ■ Compute/ETL   ■ Data",
         sz=8, color=GRAY)

    # Color dots for legend
    _shape(slide, Inches(0.5), fy + Inches(0.18), Inches(0.12), Inches(0.12),
           ORACLE_RED, MSO_SHAPE.OVAL)
    _txt(slide, Inches(0.65), fy + Inches(0.14), Inches(1), Inches(0.2),
         "Oracle (Today)", sz=8, bold=True, color=ORACLE_RED)

    _shape(slide, Inches(2.2), fy + Inches(0.18), Inches(0.12), Inches(0.12),
           MS_BLUE, MSO_SHAPE.OVAL)
    _txt(slide, Inches(2.35), fy + Inches(0.14), Inches(1.5), Inches(0.2),
         "Fabric / Power BI (Tomorrow)", sz=8, bold=True, color=MS_BLUE)

    _shape(slide, Inches(4.5), fy + Inches(0.18), Inches(0.12), Inches(0.12),
           COPILOT_BLUE, MSO_SHAPE.OVAL)
    _txt(slide, Inches(4.65), fy + Inches(0.14), Inches(2.0), Inches(0.2),
         "Copilot & Data Agents + O365", sz=8, bold=True, color=COPILOT_BLUE)

    # Security & storage labels
    _txt(slide, Inches(7.5), fy + Inches(0.14), Inches(3), Inches(0.2),
         "Security: App Roles → Workspace Roles + RLS/OLS",
         sz=8, color=GRAY)

    _txt(slide, Inches(10.5), fy + Inches(0.14), Inches(3), Inches(0.2),
         "Storage: Oracle DB → OneLake (Delta)",
         sz=8, color=GRAY)

    prs.save(output)
    return output


def main():
    parser = argparse.ArgumentParser(
        description="Generate OAC/Essbase vs Fabric architecture comparison slide")
    parser.add_argument("--output", "-o",
                        default="output/OAC_to_Fabric_Architecture.pptx")
    args = parser.parse_args()
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    path = build(args.output)
    print(f"✅ Architecture slide generated: {path}")
    print(f"   1 slide • widescreen 16:9 • {os.path.getsize(path):,} bytes")


if __name__ == "__main__":
    main()
