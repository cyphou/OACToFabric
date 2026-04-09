#!/usr/bin/env python3
"""Generate a single slide mapping OAC/Essbase services to Fabric/Power BI equivalents.

Usage:
    python scripts/generate_equivalence_slide.py [--output path/to/slide.pptx]
"""
from __future__ import annotations

import argparse
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x10, 0x7C, 0x10)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
TEAL = RGBColor(0x00, 0xB7, 0xC3)
PURPLE = RGBColor(0x6B, 0x2D, 0x8B)
ORACLE_RED = RGBColor(0xC7, 0x40, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF2, 0xF1)
GRAY = RGBColor(0x60, 0x60, 0x60)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _shape(slide, l, t, w, h, fill, st=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def _txt(slide, l, t, w, h, text, sz=18, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, name="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return tb


def _row(slide, y, oracle_label, oracle_desc, fabric_label, fabric_desc,
         accent, pct, row_bg=WHITE):
    """One mapping row: Oracle service → Fabric equivalent + conversion rate."""
    h = Inches(0.55)

    # Background
    _shape(slide, Inches(0.3), y, Inches(12.7), h, row_bg, MSO_SHAPE.RECTANGLE)

    # Oracle side
    dot = _shape(slide, Inches(0.45), y + Inches(0.17),
                 Inches(0.14), Inches(0.14), ORACLE_RED, MSO_SHAPE.OVAL)
    _txt(slide, Inches(0.7), y + Inches(0.05), Inches(2.5), Inches(0.25),
         oracle_label, sz=11, bold=True, color=ORACLE_RED)
    _txt(slide, Inches(0.7), y + Inches(0.28), Inches(3.5), Inches(0.22),
         oracle_desc, sz=8.5, color=GRAY)

    # Arrow
    _txt(slide, Inches(4.55), y + Inches(0.08), Inches(0.5), Inches(0.35),
         "→", sz=18, bold=True, color=accent, align=PP_ALIGN.CENTER)

    # Fabric side
    dot2 = _shape(slide, Inches(5.15), y + Inches(0.17),
                  Inches(0.14), Inches(0.14), accent, MSO_SHAPE.OVAL)
    _txt(slide, Inches(5.4), y + Inches(0.05), Inches(2.8), Inches(0.25),
         fabric_label, sz=11, bold=True, color=accent)
    _txt(slide, Inches(5.4), y + Inches(0.28), Inches(4.8), Inches(0.22),
         fabric_desc, sz=8.5, color=GRAY)

    # Conversion rate bar + label
    bar_left = Inches(10.6)
    bar_w = Inches(1.8)
    bar_h = Inches(0.18)
    # Background bar
    _shape(slide, bar_left, y + Inches(0.09), bar_w, bar_h,
           LIGHT, MSO_SHAPE.ROUNDED_RECTANGLE)
    # Fill bar
    pct_color = GREEN if pct >= 90 else (ORANGE if pct >= 60 else ORACLE_RED)
    fill_w = int(bar_w * pct / 100)
    if fill_w > 0:
        _shape(slide, bar_left, y + Inches(0.09), fill_w, bar_h,
               pct_color, MSO_SHAPE.ROUNDED_RECTANGLE)
    # Percentage text
    _txt(slide, bar_left, y + Inches(0.28), bar_w, Inches(0.22),
         f"{pct}% automated", sz=8, bold=True, color=pct_color,
         align=PP_ALIGN.CENTER)


def build(output: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)

    # ══ Title ══
    _shape(slide, 0, 0, SLIDE_W, Inches(1.0), DARK_BLUE, MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(0.6), Inches(0.15), Inches(11), Inches(0.45),
         "Oracle Analytics & Essbase  →  Microsoft Fabric & Power BI",
         sz=24, bold=True, color=WHITE)
    _txt(slide, Inches(0.6), Inches(0.55), Inches(10), Inches(0.35),
         "Service-by-service equivalence with conversion rates",
         sz=13, color=TEAL)

    # ══ Column headers ══
    y0 = Inches(1.15)
    _shape(slide, Inches(0.3), y0, Inches(4.6), Inches(0.4),
           RGBColor(0xFD, 0xEE, 0xE8), MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(0.5), y0 + Inches(0.05), Inches(4.4), Inches(0.3),
         "Oracle (Today)", sz=13, bold=True, color=ORACLE_RED,
         align=PP_ALIGN.CENTER)

    _shape(slide, Inches(5.5), y0, Inches(5.0), Inches(0.4),
           RGBColor(0xE5, 0xF0, 0xFD), MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(5.5), y0 + Inches(0.05), Inches(5.0), Inches(0.3),
         "Microsoft Fabric & Power BI (Tomorrow)", sz=13, bold=True,
         color=MS_BLUE, align=PP_ALIGN.CENTER)

    _shape(slide, Inches(10.5), y0, Inches(2.5), Inches(0.4),
           RGBColor(0xE5, 0xFD, 0xE8), MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(10.5), y0 + Inches(0.05), Inches(2.5), Inches(0.3),
         "Conversion Rate", sz=13, bold=True,
         color=GREEN, align=PP_ALIGN.CENTER)

    # ══ Mapping rows ══
    rows = [
        # (oracle_label, oracle_desc, fabric_label, fabric_desc, accent, pct)
        ("Oracle Database / Exadata",
         "Relational database, tables, views, stored procedures",
         "Fabric Lakehouse / Warehouse",
         "Delta tables, SQL analytics endpoint, T-SQL views",
         MS_BLUE, 95),

        ("Oracle Analytics Cloud (OAC)",
         "Enterprise BI platform — reports, dashboards, answers",
         "Power BI Service",
         "Reports, dashboards, apps, subscriptions, embedding",
         MS_BLUE, 97),

        ("RPD Semantic Model",
         "Three-layer data model (physical / logical / presentation)",
         "Power BI Semantic Model (TMDL)",
         "Tables, relationships, measures, hierarchies, perspectives",
         TEAL, 95),

        ("OAC Analyses & Dashboards",
         "Interactive reports with 80+ visual types, prompts, actions",
         "Power BI Reports (PBIR)",
         "80+ visual types, slicers, bookmarks, drillthrough, Q&A",
         TEAL, 85),

        ("OAC Data Flows",
         "ETL pipelines — source, transform, load to data model",
         "Fabric Data Pipelines + Dataflow Gen2",
         "Copy activities, transformations, PySpark notebooks, schedules",
         GREEN, 62),

        ("Essbase Cubes (BSO / ASO)",
         "Multi-dimensional OLAP — dimensions, hierarchies, aggregations",
         "Power BI Semantic Model (in-memory)",
         "Star schema, DAX measures, hierarchies, calculation groups",
         ORANGE, 78),

        ("Essbase Calc Scripts & MDX",
         "@SUM, @PRIOR, FIX/ENDFIX, YTD, business rules, allocations",
         "DAX Measures & Calculated Tables",
         "CALCULATE, DATEADD, TOTALYTD, SWITCH, Time Intelligence",
         ORANGE, 65),

        ("Smart View (Excel Add-in)",
         "HsGetValue, ad-hoc grids, POV bar, member selection, data forms",
         "Excel + CUBE Functions / Analyze in Excel",
         "CUBEVALUE, PivotTables, slicers, XMLA endpoint, Power BI add-in",
         PURPLE, 70),

        ("BI Publisher",
         "Pixel-perfect operational reports (invoices, statements, labels)",
         "Power BI Paginated Reports (.rdl)",
         "Pixel-perfect layouts, parameters, export to PDF/Excel/Word",
         PURPLE, 75),

        ("OAC Security (Roles / Session Vars)",
         "Application roles, session variable filters, object permissions",
         "Fabric Workspace Roles + RLS / OLS",
         "Workspace admin/member/viewer, DAX row-level & object-level security",
         GREEN, 90),

        ("OAC Agents & Alerts",
         "Scheduled condition-based alerts, email notifications",
         "Data Activator + Power Automate",
         "Real-time triggers, reflex actions, Teams/email notifications",
         MS_BLUE, 50),
    ]

    y = Inches(1.7)
    row_h = Inches(0.55)
    for i, (ol, od, fl, fd, accent, pct) in enumerate(rows):
        bg = LIGHT if i % 2 == 0 else WHITE
        _row(slide, y, ol, od, fl, fd, accent, pct, bg)
        y += row_h

    prs.save(output)
    return output


def main():
    parser = argparse.ArgumentParser(
        description="Generate OAC/Essbase → Fabric service equivalence slide")
    parser.add_argument("--output", "-o",
                        default="output/OAC_to_Fabric_Equivalence.pptx")
    args = parser.parse_args()
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    path = build(args.output)
    print(f"✅ Equivalence slide generated: {path}")
    print(f"   1 slide • widescreen 16:9 • {os.path.getsize(path):,} bytes")


if __name__ == "__main__":
    main()
