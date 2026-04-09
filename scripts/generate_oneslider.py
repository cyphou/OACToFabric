#!/usr/bin/env python3
"""Generate a single executive one-pager slide for OAC/Essbase → Fabric migration.

Usage:
    python scripts/generate_oneslider.py [--output path/to/slide.pptx]
"""
from __future__ import annotations

import argparse
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colours ──
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x10, 0x7C, 0x10)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
TEAL = RGBColor(0x00, 0xB7, 0xC3)
PURPLE = RGBColor(0x6B, 0x2D, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF2, 0xF1)
GRAY = RGBColor(0x60, 0x60, 0x60)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _shape(slide, l, t, w, h, fill, st=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def _txt(slide, l, t, w, h, text, sz=18, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, name="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return tb


def _kpi(slide, l, t, w, h, val, label, accent):
    card = _shape(slide, l, t, w, h, WHITE)
    card.line.color.rgb = accent
    card.line.width = Pt(2.5)
    _txt(slide, l + Inches(0.05), t + Inches(0.08),
         w - Inches(0.1), Inches(0.45),
         val, sz=26, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _txt(slide, l + Inches(0.05), t + Inches(0.5),
         w - Inches(0.1), Inches(0.3),
         label, sz=9, color=GRAY, align=PP_ALIGN.CENTER)


def build(output: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)

    # ══ Title bar ══
    _shape(slide, 0, 0, SLIDE_W, Inches(1.15), DARK_BLUE, MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(0.6), Inches(0.15), Inches(8), Inches(0.55),
         "Oracle Analytics & Essbase → Microsoft Fabric & Power BI",
         sz=24, bold=True, color=WHITE)
    _txt(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.4),
         "AI-accelerated migration  —  97% object coverage  —  ~50% faster delivery",
         sz=13, color=TEAL)

    # Logos placeholder (right side of title)
    for i, (label, c) in enumerate([("Oracle", RGBColor(0xF8, 0x00, 0x00)),
                                     ("→", TEAL),
                                     ("Fabric", MS_BLUE)]):
        _txt(slide, Inches(10.2 + i * 1.0), Inches(0.3), Inches(0.9), Inches(0.5),
             label, sz=14, bold=True, color=WHITE if label != "→" else TEAL,
             align=PP_ALIGN.CENTER)

    # ══ ROW 1: Source → Target flow ══
    y1 = Inches(1.35)

    # Source box
    src = _shape(slide, Inches(0.4), y1, Inches(3.0), Inches(2.4),
                 RGBColor(0xFD, 0xF0, 0xE5))
    src.line.color.rgb = ORANGE
    src.line.width = Pt(1.5)
    _txt(slide, Inches(0.6), y1 + Inches(0.08), Inches(2.6), Inches(0.35),
         "TODAY", sz=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    sources = [
        "Oracle Analytics Cloud (OAC)",
        "Oracle Essbase (OLAP Cubes)",
        "Smart View (Excel Add-in)",
        "Oracle Database (Exadata)",
        "RPD Semantic Models",
        "BI Publisher Reports",
    ]
    sy = y1 + Inches(0.45)
    for s in sources:
        _txt(slide, Inches(0.65), sy, Inches(2.7), Inches(0.22),
             f"•  {s}", sz=9, color=BLACK)
        sy += Inches(0.28)

    # Arrow
    _txt(slide, Inches(3.5), y1 + Inches(0.85), Inches(0.6), Inches(0.6),
         "→", sz=36, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Target box
    tgt = _shape(slide, Inches(4.2), y1, Inches(3.0), Inches(2.4),
                 RGBColor(0xE5, 0xF0, 0xFD))
    tgt.line.color.rgb = MS_BLUE
    tgt.line.width = Pt(1.5)
    _txt(slide, Inches(4.4), y1 + Inches(0.08), Inches(2.6), Inches(0.35),
         "TOMORROW", sz=13, bold=True, color=MS_BLUE, align=PP_ALIGN.CENTER)
    targets = [
        "Microsoft Fabric (Lakehouse)",
        "Power BI Semantic Models",
        "Power BI Reports & Dashboards",
        "Excel + CUBE Functions",
        "Fabric Data Pipelines",
        "Data Activator (Alerts)",
    ]
    ty = y1 + Inches(0.45)
    for t in targets:
        _txt(slide, Inches(4.45), ty, Inches(2.7), Inches(0.22),
             f"•  {t}", sz=9, color=BLACK)
        ty += Inches(0.28)

    # ══ ROW 1 right: KPI cards ══
    kpis = [
        ("–50%", "Time to\nMigrate", GREEN),
        ("–60%", "Infra\nTCO", MS_BLUE),
        ("–50%", "Ops\nEffort", TEAL),
        ("97%", "Object\nCoverage", ORANGE),
    ]
    kx = Inches(7.7)
    for val, label, color in kpis:
        _kpi(slide, kx, y1, Inches(1.3), Inches(0.85), val, label, color)
        kx += Inches(1.4)

    # Second row of KPIs
    kpis2 = [
        ("185", "Object\nMappings", PURPLE),
        ("120+", "DAX\nRules", ORANGE),
        ("80+", "Visual\nMappings", TEAL),
        ("4,005", "Automated\nTests", GREEN),
    ]
    kx = Inches(7.7)
    for val, label, color in kpis2:
        _kpi(slide, kx, y1 + Inches(1.05), Inches(1.3), Inches(0.85),
             val, label, color)
        kx += Inches(1.4)

    # Payback
    _shape(slide, Inches(7.7), y1 + Inches(2.1), Inches(5.3), Inches(0.3),
           GREEN, MSO_SHAPE.ROUNDED_RECTANGLE)
    _txt(slide, Inches(7.7), y1 + Inches(2.1), Inches(5.3), Inches(0.3),
         "Payback: 9–14 months  •  3-Year Net Savings: $2M–$6M+",
         sz=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ══ ROW 2: What migrates ══
    y2 = Inches(4.05)

    _txt(slide, Inches(0.4), y2, Inches(3.5), Inches(0.3),
         "What Migrates Automatically", sz=14, bold=True, color=DARK_BLUE)

    cols = [
        ("Data Layer", [
            "Tables & views → Delta tables",
            "25+ data type mappings",
            "DDL auto-generated",
        ], MS_BLUE),
        ("Analytics Layer", [
            "RPD model → TMDL semantic model",
            "OAC expressions → DAX measures",
            "Essbase cubes → semantic models",
        ], TEAL),
        ("Reporting Layer", [
            "80+ visual types → Power BI",
            "Prompts → slicers & parameters",
            "Smart View → Excel CUBE",
        ], ORANGE),
        ("Security Layer", [
            "App roles → workspace roles",
            "Row-level & object-level security",
            "Session vars → DAX RLS",
        ], PURPLE),
    ]
    cx = Inches(0.4)
    cw = Inches(3.05)
    for title, items, color in cols:
        bar = _shape(slide, cx, y2 + Inches(0.35), cw, Inches(0.28),
                     color, MSO_SHAPE.ROUNDED_RECTANGLE)
        bar.text_frame.paragraphs[0].text = f"  {title}"
        bar.text_frame.paragraphs[0].font.size = Pt(10)
        bar.text_frame.paragraphs[0].font.bold = True
        bar.text_frame.paragraphs[0].font.color.rgb = WHITE
        bar.text_frame.paragraphs[0].font.name = "Segoe UI"
        iy = y2 + Inches(0.68)
        for item in items:
            _txt(slide, cx + Inches(0.1), iy, cw - Inches(0.15), Inches(0.2),
                 f"✓  {item}", sz=8.5, color=BLACK)
            iy += Inches(0.2)
        cx += cw + Inches(0.15)

    # ══ ROW 3: Timeline + Next Steps ══
    y3 = Inches(5.65)

    _shape(slide, 0, y3, SLIDE_W, Inches(1.85), RGBColor(0xF7, 0xF7, 0xF7),
           MSO_SHAPE.RECTANGLE)

    _txt(slide, Inches(0.5), y3 + Inches(0.1), Inches(3), Inches(0.3),
         "Typical Timeline", sz=13, bold=True, color=DARK_BLUE)

    phases = [
        ("Month 1–2", "Discovery &\nSchema", MS_BLUE),
        ("Month 3–4", "ETL & Semantic\nModels", GREEN),
        ("Month 5–6", "Reports &\nSecurity", ORANGE),
        ("Month 7+", "Validation &\nGo-Live", TEAL),
    ]
    px = Inches(0.5)
    pw = Inches(1.7)
    for period, desc, color in phases:
        dot = _shape(slide, px + Inches(0.7), y3 + Inches(0.45),
                     Inches(0.18), Inches(0.18), color, MSO_SHAPE.OVAL)
        if period != "Month 7+":
            _shape(slide, px + pw, y3 + Inches(0.5),
                   Inches(0.45), Inches(0.04), LIGHT, MSO_SHAPE.RECTANGLE)
        _txt(slide, px + Inches(0.15), y3 + Inches(0.65), pw, Inches(0.22),
             period, sz=10, bold=True, color=color, align=PP_ALIGN.CENTER)
        _txt(slide, px + Inches(0.15), y3 + Inches(0.85), pw, Inches(0.4),
             desc, sz=8.5, color=GRAY, align=PP_ALIGN.CENTER)
        px += pw + Inches(0.45)

    # Next Steps
    _txt(slide, Inches(8.0), y3 + Inches(0.1), Inches(5), Inches(0.3),
         "Recommended Next Steps", sz=13, bold=True, color=DARK_BLUE)

    steps = [
        ("1", "Discovery Workshop — automated inventory of your OAC/Essbase estate"),
        ("2", "Proof of Concept — migrate a representative subset to validate fidelity"),
        ("3", "Wave Planning — phased rollout with automated validation & rollback"),
    ]
    stepy = y3 + Inches(0.45)
    for num, desc in steps:
        circle = _shape(slide, Inches(8.1), stepy, Inches(0.25), Inches(0.25),
                        MS_BLUE, MSO_SHAPE.OVAL)
        circle.text_frame.paragraphs[0].text = num
        circle.text_frame.paragraphs[0].font.size = Pt(9)
        circle.text_frame.paragraphs[0].font.bold = True
        circle.text_frame.paragraphs[0].font.color.rgb = WHITE
        circle.text_frame.paragraphs[0].font.name = "Segoe UI"
        circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        _txt(slide, Inches(8.5), stepy + Inches(0.02), Inches(4.5), Inches(0.22),
             desc, sz=9, color=BLACK)
        stepy += Inches(0.35)

    prs.save(output)
    return output


def main():
    parser = argparse.ArgumentParser(description="Generate OAC→Fabric executive one-slider")
    parser.add_argument("--output", "-o",
                        default="output/OAC_to_Fabric_OnePager.pptx")
    args = parser.parse_args()
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    path = build(args.output)
    print(f"✅ One-slider generated: {path}")
    print(f"   1 slide • widescreen 16:9 • {os.path.getsize(path):,} bytes")


if __name__ == "__main__":
    main()
