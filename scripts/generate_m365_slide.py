#!/usr/bin/env python3
"""Generate a slide comparing M365 integration: OAC vs Fabric/Power BI.

Usage:
    python scripts/generate_m365_slide.py [--output path/to/slide.pptx]
"""
from __future__ import annotations

import argparse
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette ──
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
GREEN = RGBColor(0x10, 0x7C, 0x10)
TEAL = RGBColor(0x00, 0xB7, 0xC3)
PURPLE = RGBColor(0x6B, 0x2D, 0x8B)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ORACLE_RED = RGBColor(0xC7, 0x40, 0x34)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF2, 0xF1)
LIGHT_RED = RGBColor(0xFD, 0xEE, 0xE8)
LIGHT_BLUE = RGBColor(0xE5, 0xF0, 0xFD)
LIGHT_GREEN = RGBColor(0xE5, 0xFD, 0xE8)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
BLACK = RGBColor(0x00, 0x00, 0x00)
COPILOT = RGBColor(0x6C, 0x2B, 0xD9)

# M365 app brand colors
EXCEL_GREEN = RGBColor(0x21, 0x7A, 0x46)
PPT_ORANGE = RGBColor(0xD2, 0x4B, 0x27)
WORD_BLUE = RGBColor(0x29, 0x5A, 0xB8)
TEAMS_PURPLE = RGBColor(0x50, 0x59, 0xC9)
OUTLOOK_BLUE = RGBColor(0x00, 0x78, 0xD4)
SP_TEAL = RGBColor(0x03, 0x82, 0x7F)
OD_BLUE = RGBColor(0x09, 0x4A, 0xB2)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Logo directory (resolved relative to this script)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_DIR = os.path.join(os.path.dirname(SCRIPT_DIR), "assets", "logos")

# Map app display names to logo filenames
LOGO_FILES = {
    "Excel": "excel.png",
    "PowerPoint": "powerpoint.png",
    "Word": "word.png",
    "Teams": "teams.png",
    "Outlook": "outlook.png",
    "SharePoint": "sharepoint.png",
    "OneDrive": "onedrive.png",
    "Copilot M365": "copilot.png",
}


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _shape(slide, l, t, w, h, fill, st=MSO_SHAPE.ROUNDED_RECTANGLE,
           border=None, border_w=Pt(1)):
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.fill.solid()
        s.line.fill.fore_color.rgb = border
        s.line.width = border_w
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s


def _txt(slide, l, t, w, h, text, sz=18, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, name="Segoe UI"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    return tb


def _multi_txt(slide, l, t, w, h, lines, sz=9, color=BLACK, name="Segoe UI"):
    """Textbox with multiple lines."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, clr, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = clr
        p.font.name = name
        p.space_after = Pt(2)
    return tb


def _icon_badge(slide, l, t, color, label, sz=9):
    """Small colored badge with label."""
    w = Inches(0.18)
    _shape(slide, l, t, w, w, color, MSO_SHAPE.ROUNDED_RECTANGLE)
    _txt(slide, l + Inches(0.22), t - Inches(0.01), Inches(1.2), Inches(0.22),
         label, sz=sz, bold=True, color=color)


def _check(slide, l, t, color=GREEN):
    """Green checkmark."""
    _txt(slide, l, t, Inches(0.2), Inches(0.2), "✓", sz=11, bold=True,
         color=color, align=PP_ALIGN.CENTER)


def _cross(slide, l, t, color=ORACLE_RED):
    """Red cross."""
    _txt(slide, l, t, Inches(0.2), Inches(0.2), "✗", sz=11, bold=True,
         color=color, align=PP_ALIGN.CENTER)


def _partial(slide, l, t, color=ORANGE):
    """Orange tilde for partial support."""
    _txt(slide, l, t, Inches(0.2), Inches(0.2), "~", sz=13, bold=True,
         color=color, align=PP_ALIGN.CENTER)


def build(output: str) -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)

    # ══ Title bar ══
    _shape(slide, 0, 0, SLIDE_W, Inches(0.85), DARK_BLUE, st=MSO_SHAPE.RECTANGLE)
    _txt(slide, Inches(0.5), Inches(0.1), Inches(12), Inches(0.4),
         "Microsoft 365 Integration:  Oracle OAC  vs  Microsoft Fabric",
         sz=22, bold=True, color=WHITE)
    _txt(slide, Inches(0.5), Inches(0.48), Inches(10), Inches(0.3),
         "Fabric is natively embedded in M365 — OAC requires manual bridging",
         sz=11, color=TEAL)

    # ══ Column headers ══
    y0 = Inches(1.05)
    col_app = Inches(0.3)
    col_oac = Inches(3.5)
    col_fabric = Inches(7.8)
    col_w_app = Inches(3.0)
    col_w_oac = Inches(4.1)
    col_w_fabric = Inches(5.2)

    # App column header
    _shape(slide, col_app, y0, col_w_app, Inches(0.38),
           RGBColor(0xE8, 0xE8, 0xE8), MSO_SHAPE.RECTANGLE)
    _txt(slide, col_app + Inches(0.1), y0 + Inches(0.05), col_w_app, Inches(0.28),
         "M365 Application", sz=12, bold=True, color=DARK_BLUE,
         align=PP_ALIGN.CENTER)

    # OAC column header
    _shape(slide, col_oac, y0, col_w_oac, Inches(0.38),
           LIGHT_RED, MSO_SHAPE.RECTANGLE)
    _txt(slide, col_oac, y0 + Inches(0.05), col_w_oac, Inches(0.28),
         "Oracle Analytics Cloud", sz=12, bold=True, color=ORACLE_RED,
         align=PP_ALIGN.CENTER)

    # Fabric column header
    _shape(slide, col_fabric, y0, col_w_fabric, Inches(0.38),
           LIGHT_BLUE, MSO_SHAPE.RECTANGLE)
    _txt(slide, col_fabric, y0 + Inches(0.05), col_w_fabric, Inches(0.28),
         "Microsoft Fabric & Power BI", sz=12, bold=True, color=MS_BLUE,
         align=PP_ALIGN.CENTER)

    # ══ M365 app rows ══
    apps = [
        # (app_name, brand_color, oac_icon, oac_lines, fabric_icon, fabric_lines)
        ("Excel", EXCEL_GREEN,
         "partial",
         [("Smart View add-in (proprietary)", GRAY, False),
          ("Requires Oracle client install", GRAY, False),
          ("No live data refresh", ORACLE_RED, False)],
         "check",
         [("Analyze in Excel — live PivotTables", GRAY, False),
          ("CUBE functions (CUBEVALUE, CUBEMEMBER)", GRAY, False),
          ("XMLA endpoint, Power BI add-in", GREEN, False),
          ("Copilot in Excel pulls from semantic model", COPILOT, False)]),

        ("PowerPoint", PPT_ORANGE,
         "cross",
         [("Manual screenshot + paste", GRAY, False),
          ("No live connection", ORACLE_RED, False),
          ("Static images only", GRAY, False)],
         "check",
         [("Live Power BI visuals in slides", GRAY, False),
          ("Auto-refresh on presentation open", GRAY, False),
          ("Storytelling with data add-in", GREEN, False),
          ("Copilot generates slides from data", COPILOT, False)]),

        ("Word", WORD_BLUE,
         "cross",
         [("No integration", ORACLE_RED, False),
          ("Manual export → copy/paste", GRAY, False)],
         "check",
         [("Copilot pulls data from semantic models", GRAY, False),
          ("Insert live charts & tables", GRAY, False),
          ("Auto-generated data narratives", GREEN, False)]),

        ("Teams", TEAMS_PURPLE,
         "cross",
         [("No native embedding", ORACLE_RED, False),
          ("Requires custom iFrame / URL links", GRAY, False),
          ("No interactive experience", GRAY, False)],
         "check",
         [("Native Power BI tab in channels", GRAY, False),
          ("Chat with data (Q&A in Teams)", GRAY, False),
          ("Data Agents accessible in Teams", GREEN, False),
          ("Meeting-integrated analytics", GREEN, False)]),

        ("Outlook", OUTLOOK_BLUE,
         "partial",
         [("OAC Agents send basic email alerts", GRAY, False),
          ("Plain-text notifications only", GRAY, False),
          ("No rich cards or actions", ORACLE_RED, False)],
         "check",
         [("Interactive report cards in email", GRAY, False),
          ("Data Activator alerts with rich actions", GRAY, False),
          ("Subscribe to reports → inbox delivery", GREEN, False),
          ("Power Automate → actionable notifications", GREEN, False)]),

        ("SharePoint", SP_TEAL,
         "cross",
         [("No native integration", ORACLE_RED, False),
          ("Manual iFrame embed only", GRAY, False)],
         "check",
         [("Power BI web parts for SharePoint pages", GRAY, False),
          ("Embedded interactive reports", GRAY, False),
          ("Auto page creation from datasets", GREEN, False)]),

        ("OneDrive", OD_BLUE,
         "cross",
         [("No integration", ORACLE_RED, False),
          ("Files stored in Oracle Cloud only", GRAY, False)],
         "check",
         [("Shared .pbip project files", GRAY, False),
          ("Collaborative editing & version control", GRAY, False),
          ("Git integration for semantic models", GREEN, False)]),

        ("Copilot M365", COPILOT,
         "cross",
         [("No AI assistant integration", ORACLE_RED, False),
          ("No natural language data access", GRAY, False)],
         "check",
         [("Copilot queries Fabric semantic models", GRAY, False),
          ("Data Agents answer in Teams, Excel, Word", COPILOT, False),
          ("Natural language \u2192 DAX / SQL / visuals", COPILOT, False),
          ("Grounded in enterprise data governance", GREEN, False)]),
    ]

    y = Inches(1.55)
    row_h = Inches(0.68)

    for i, (app_name, app_color, oac_icon, oac_lines,
            fab_icon, fab_lines) in enumerate(apps):
        bg = LIGHT if i % 2 == 0 else WHITE
        # Row background
        _shape(slide, Inches(0.3), y, Inches(12.7), row_h,
               bg, MSO_SHAPE.RECTANGLE)

        # App logo + name
        logo_file = LOGO_FILES.get(app_name)
        logo_path = os.path.join(LOGO_DIR, logo_file) if logo_file else None
        if logo_path and os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(0.38), y + Inches(0.08),
                Inches(0.48), Inches(0.48))
        else:
            _shape(slide, Inches(0.42), y + Inches(0.12),
                   Inches(0.42), Inches(0.42), app_color,
                   MSO_SHAPE.ROUNDED_RECTANGLE)
            _txt(slide, Inches(0.42), y + Inches(0.15), Inches(0.42),
                 Inches(0.36), app_name[0], sz=16, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        _txt(slide, Inches(0.95), y + Inches(0.18), Inches(2.2), Inches(0.28),
             app_name, sz=12, bold=True, color=app_color)

        # OAC status icon
        if oac_icon == "check":
            _check(slide, col_oac + Inches(0.05), y + Inches(0.08))
        elif oac_icon == "cross":
            _cross(slide, col_oac + Inches(0.05), y + Inches(0.08))
        else:
            _partial(slide, col_oac + Inches(0.05), y + Inches(0.05))

        # OAC details
        _multi_txt(slide, col_oac + Inches(0.3), y + Inches(0.06),
                   col_w_oac - Inches(0.4), row_h - Inches(0.1),
                   oac_lines, sz=8)

        # Fabric status icon
        _check(slide, col_fabric + Inches(0.05), y + Inches(0.08))

        # Fabric details
        _multi_txt(slide, col_fabric + Inches(0.3), y + Inches(0.06),
                   col_w_fabric - Inches(0.4), row_h - Inches(0.1),
                   fab_lines, sz=8)

        y += row_h

    # ══ Footer / legend ══
    fy = Inches(7.08)
    _shape(slide, Inches(0.3), fy, Inches(12.7), Inches(0.35),
           RGBColor(0xF8, 0xF8, 0xF8), MSO_SHAPE.RECTANGLE)

    _check(slide, Inches(0.5), fy + Inches(0.05))
    _txt(slide, Inches(0.7), fy + Inches(0.07), Inches(1.2), Inches(0.2),
         "Native integration", sz=8, color=GREEN)

    _partial(slide, Inches(2.1), fy + Inches(0.02))
    _txt(slide, Inches(2.3), fy + Inches(0.07), Inches(1.2), Inches(0.2),
         "Limited / manual", sz=8, color=ORANGE)

    _cross(slide, Inches(3.8), fy + Inches(0.05))
    _txt(slide, Inches(4.0), fy + Inches(0.07), Inches(1.2), Inches(0.2),
         "Not available", sz=8, color=ORACLE_RED)

    _shape(slide, Inches(5.6), fy + Inches(0.1), Inches(0.14), Inches(0.14),
           COPILOT, MSO_SHAPE.OVAL)
    _txt(slide, Inches(5.8), fy + Inches(0.07), Inches(2), Inches(0.2),
         "Copilot / Data Agent powered", sz=8, color=COPILOT)

    _txt(slide, Inches(8.5), fy + Inches(0.07), Inches(4.5), Inches(0.2),
         "Fabric is part of M365 — OAC sits outside the Microsoft ecosystem",
         sz=8, bold=True, color=DARK_BLUE)

    prs.save(output)
    return output


def main():
    parser = argparse.ArgumentParser(
        description="Generate M365 integration comparison slide")
    parser.add_argument("--output", "-o",
                        default="output/OAC_to_Fabric_M365_Integration.pptx")
    args = parser.parse_args()
    os.makedirs(os.path.dirname(args.output) or ".", exist_ok=True)
    path = build(args.output)
    print(f"✅ M365 integration slide generated: {path}")
    print(f"   1 slide • widescreen 16:9 • {os.path.getsize(path):,} bytes")


if __name__ == "__main__":
    main()
